#!/usr/bin/env python3
from __future__ import annotations
from collections import Counter
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile
import json, math, re, shutil, sys
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

ROOT=Path(__file__).resolve().parents[1]
sys.path.insert(0,str(ROOT/'tools'))
import build_b394_delivery as base
import package_youtube_delivery as youtube

ID='B396-SET-001'; PROJECT='前值转移三码'; PERIODS=30; P0=.30
INPUT=ROOT/'01_本次输入'/'哈希分分彩_20260731_0181至0380.txt'
OUT=ROOT/'dist'/f'{PROJECT}_交付'; PKG=OUT/'package'
SCHEME=OUT/f'{PROJECT}_方案套.zip'; PPT=OUT/f'{PROJECT}_挂机前讲解.pptx'
EVIDENCE=OUT/f'{PROJECT}_验证记录.json'; MANIFEST=OUT/'DELIVERY_MANIFEST.json'
SEO=ROOT/'youtube_seo'/f'{PROJECT}.json'; POS=[(2,'百位'),(3,'十位'),(4,'个位')]


def row_counts(values,state):
 c=Counter()
 for a,b in zip(values[:-1],values[1:]):
  if a==state:c[b]+=1
 return c

def select(values,state):
 row=row_counts(values,state); overall=Counter(values[1:]); support=sum(row.values()); fallback=support<5
 rank=sorted(range(10),key=(lambda d:(-overall[d],d)) if fallback else (lambda d:(-row[d],-overall[d],d)))
 return rank[:3],row,fallback

def binom_tail(n,h,p=P0):
 return min(1.0,sum(math.comb(n,k)*p**k*(1-p)**(n-k) for k in range(h,n+1)))
def miss_streak(xs):
 best=run=0
 for x in xs:
  run=0 if x else run+1; best=max(best,run)
 return best

def summary(xs):
 n=len(xs); h=sum(xs)
 return {'预测次数':n,'命中次数':h,'命中比例':round(h/n,6),'随机三码基准':P0,'相对基准差':round(h/n-P0,6),'独立近似单侧二项P值':round(binom_tail(n,h),6),'最大连续未中':miss_streak(xs)}
def audit(rows):
 seg={'校准段':[],'验证段':[],'审计段':[]}; per={n:[] for _,n in POS}; allx=[]
 for t in range(60,len(rows)):
  for i,n in POS:
   hist=[x[i] for _,x in rows[:t]]; pred,_,_=select(hist,hist[-1]); hit=rows[t][1][i] in pred
   allx.append(hit); per[n].append(hit); seg['校准段' if t<120 else '验证段' if t<160 else '审计段'].append(hit)
 return {'方法':'从第61期开始逐期扩展窗口；每次只使用此前数据建立同位置一阶转移频次。','总计':summary(allx),'分段':{k:summary(v) for k,v in seg.items()},'分位置':{k:summary(v) for k,v in per.items()},'统计边界':'二项P值仅作独立近似参考；位置间和时间上可能相关，不作为显著性结论。','解释边界':'复用数据审计，不属于独立样本外；最终结论必须来自新的实际挂机记录。'}
def freeze(rows):
 out={}
 for i,n in POS:
  values=[x[i] for _,x in rows]; state=values[-1]; top,row,fallback=select(values,state)
  overall=Counter(values[1:]); ranked=[(d,row[d]) for d in sorted(range(10),key=(lambda d:(-overall[d],d)) if fallback else (lambda d:(-row[d],-overall[d],d)))]
  out[n]={'位置索引':i,'最新前值':state,'前值样本次数':sum(row.values()),'后继原始计数':{str(d):row[d] for d in range(10)},'后继计数排序':[{'数字':d,'次数':c} for d,c in ranked],'选取规则':'后继次数降序；并列按全局后继频次降序、数字升序；取前3名。','低样本回退':fallback,'冻结号码':top,'运行中是否更新':False}
 return out


def fixed_txt(play,digits):
 lines=base.common('定码轮换',play,True)+['换号规则=9',f'换号期数={PERIODS}']+base.tail()+[f"定码轮换内容={' '.join(map(str,digits))}",'定码轮换单组=True','SchemeCreator=']
 return '\r\n'.join(lines)+'\r\n'
def build_package(g,issue,data_range):
 if PKG.exists():shutil.rmtree(PKG)
 main=PKG/'01_主方案_三位置轮投'; ctrl=PKG/'02_随机对照_单独运行'; main.mkdir(parents=True); ctrl.mkdir(parents=True)
 for n in ['百位','十位','个位']:
  (main/f'前值转移{n}-定码轮换.txt').write_bytes(fixed_txt(n,g[n]['冻结号码']).encode('gbk'))
  (ctrl/f'随机{n}-随机出号.txt').write_bytes(base.random_txt(n).encode('gbk'))
 readme=f'''# {PROJECT}｜挂机前说明

研究同一位置的一阶转移：当上一期某位置出现数字X时，历史上下一期最常跟随X的三个数字。

## 冻结号码
- 百位：前值{g['百位']['最新前值']} → {' '.join(map(str,g['百位']['冻结号码']))}
- 十位：前值{g['十位']['最新前值']} → {' '.join(map(str,g['十位']['冻结号码']))}
- 个位：前值{g['个位']['最新前值']} → {' '.join(map(str,g['个位']['冻结号码']))}
- 数据：{data_range}
- 截止期：{issue}

号码挂机前计算并冻结；软件只执行定码，不会自动重算转移频次。

## 运行
1. 只导入主方案文件夹三份TXT，手工开启顶部“方案轮投”，运行{PERIODS}期。
2. 主方案停止后，再单独运行随机对照{PERIODS}期，禁止并投。
3. 每期3个号码、每个1元；每阶段{PERIODS*3}元，建议准备{PERIODS*6}元。
4. 止盈不设置，止损不设置；每阶段满{PERIODS}期硬停止；全程平倍，禁止追损和改码。
'''
 (PKG/'00_使用说明.md').write_text(readme,encoding='utf-8')
 (PKG/'01_导入运行核对表.md').write_text(f'''# 导入与运行核对表
- [ ] 主方案与随机对照未同时启用
- [ ] 顶部“方案轮投”已开启
- [ ] 每期只运行一个位置方案
- [ ] 每期3元，全程平倍
- [ ] 每阶段满{PERIODS}期停止
- [ ] 已记录位置、冻结号码、开奖号和是否命中
''',encoding='utf-8')
 (PKG/'02_实际挂机记录模板.csv').write_text('阶段,期号,实际方案名称,位置,冻结号码,开奖号,是否命中,备注\n',encoding='utf-8-sig')
 with ZipFile(SCHEME,'w',ZIP_DEFLATED) as z:
  for p in sorted(PKG.rglob('*')):
   if p.is_file():z.write(p,arcname=p.relative_to(PKG).as_posix())


def counts(g,n,limit=5,multiline=False):
 items=[f"{i}. {x['数字']}（{x['次数']}次）" for i,x in enumerate(g[n]['后继计数排序'][:limit],1)]
 return '\n'.join(items) if multiline else '  '.join(items)
def build_ppt(g,issue):
 cover=ROOT/'assets'/'ppt'/'fixed_pages'/'首页背景图谱.png'; end=ROOT/'assets'/'ppt'/'fixed_pages'/'固定最后一页_画面.png'; prs=base.fixed.new_prs()
 s=base.fixed.add_cover(prs,cover); base.fixed.add_text(s,Inches(.82),Inches(1.02),Inches(8.8),Inches(.7),PROJECT,34,base.COLORS['white'],True); base.fixed.add_text(s,Inches(.84),Inches(1.84),Inches(9.4),Inches(.38),'上一位数字之后，下一位最常跟谁？',17,base.COLORS['gold'],True); base.fixed.add_text(s,Inches(.84),Inches(5.88),Inches(9.6),Inches(.34),'挂机前规则说明｜结果等待实际运行',14,base.COLORS['white']); base.fixed.set_notes(s,'本期研究同一位置的一阶转移频次，只讲号码来源和运行规则。')
 s=base.body_slide(prs,'这次验证什么','研究问题'); base.card(s,Inches(.78),Inches(1.8),Inches(5.55),Inches(3.8),'一个简单问题','如果百位上一期是3，历史上下一期百位最常出现哪些数字？十位和个位也分别计算。','gold',18,19); base.card(s,Inches(6.58),Inches(1.8),Inches(5.55),Inches(3.8),'不把频率当预言','转移次数只是历史条件频率，可能没有优势。主方案必须与同成本随机三码分开运行。','blue',18,18); base.fixed.set_notes(s,'解释研究问题和可证伪边界。')
 s=base.body_slide(prs,'号码怎样一步步算出来','计算方法')
 for i,(t,b) in enumerate([('按位置拆开','百位只看百位，十位只看十位，个位只看个位。'),('锁定最新前值','读取截止期最后一期该位置数字。'),('统计后继数字','历史中找出同样前值，统计下一期各数字次数。'),('取前三并冻结','按次数排序取3个；并列规则固定。')]):base.card(s,Inches(.72+(i%2)*6.12),Inches(1.78+(i//2)*2.12),Inches(5.75),Inches(1.78),f'{i+1}  {t}',b,'gold' if i%2==0 else 'blue',17,16)
 base.fixed.add_text(s,Inches(.9),Inches(6.18),Inches(11.5),Inches(.32),'号码只在挂机前计算一次，运行期间不自动更新。',17,base.COLORS['green'],True,PP_ALIGN.CENTER); base.fixed.set_notes(s,'按四步解释，观众可以复算。')
 s=base.body_slide(prs,'三组号码的真实来源','号码证据'); ac={'百位':'gold','十位':'blue','个位':'green'}
 for i,n in enumerate(['百位','十位','个位']):
  x=g[n]; body=f"最新前值：{x['最新前值']}\n历史转移样本：{x['前值样本次数']}次\n排序前三：\n{counts(g,n,3,True)}\n\n冻结号码：{' '.join(map(str,x['冻结号码']))}"; base.card(s,Inches(.72+i*4.05),Inches(1.78),Inches(3.82),Inches(4.35),n,body,ac[n],19,16)
 base.fixed.add_text(s,Inches(.88),Inches(6.3),Inches(11.4),Inches(.28),f'数据截止：{issue}｜三组数字与TXT一致',15,base.COLORS['gray'],True,PP_ALIGN.CENTER); base.fixed.set_notes(s,'读出前值、样本次数、按并列规则确定的排名和冻结号码。')
 ex=max(['百位','十位','个位'],key=lambda n:g[n]['前值样本次数']); x=g[ex]; s=base.body_slide(prs,f'以{ex}为例，完整复算一次','完整案例'); base.card(s,Inches(.78),Inches(1.82),Inches(3.55),Inches(3.9),'先找前值',f"截止期{ex}是{x['最新前值']}。\n历史中作为前值共{x['前值样本次数']}次。",'gold',18,19); base.card(s,Inches(4.58),Inches(1.82),Inches(3.55),Inches(3.9),'再数下一期',f'完整排序：\n{counts(g,ex,10)}','blue',18,16); base.card(s,Inches(8.38),Inches(1.82),Inches(3.55),Inches(3.9),'最后冻结',f"取前三：\n\n{' '.join(map(str,x['冻结号码']))}\n\n连续{PERIODS}期不改码。",'green',18,21); base.fixed.set_notes(s,'完整展示一个位置的复算过程，包括并列处理后的排序。')
 s=base.body_slide(prs,'主方案怎样运行','执行规则')
 for i,n in enumerate(['百位','十位','个位']):base.card(s,Inches(.82+i*4.08),Inches(1.95),Inches(3.52),Inches(2.02),f'{i+1}  {n}',f"冻结：{' '.join(map(str,g[n]['冻结号码']))}",ac[n],18,22)
 base.card(s,Inches(.82),Inches(4.35),Inches(11.2),Inches(1.42),'软件设置',f'导入三份主方案 → 手工开启顶部“方案轮投” → 每期一个位置 → 连续{PERIODS}期停止。','gold',17,18); base.fixed.set_notes(s,'软件只轮流执行静态定码。')
 s=base.body_slide(prs,'随机对照必须单独运行','对照设计'); base.card(s,Inches(.78),Inches(1.82),Inches(5.55),Inches(3.82),'第一阶段：主方案',f'三份前值转移定码轮投\n连续{PERIODS}期\n每期3元，共{PERIODS*3}元','gold',18,19); base.card(s,Inches(6.58),Inches(1.82),Inches(5.55),Inches(3.82),'第二阶段：随机对照',f'关闭主方案后再启用随机三码\n连续{PERIODS}期\n每期3元，共{PERIODS*3}元','blue',18,18); base.fixed.add_text(s,Inches(.9),Inches(6.02),Inches(11.5),Inches(.42),'两组不能并投，否则无法比较。',18,base.COLORS['red'],True,PP_ALIGN.CENTER); base.fixed.set_notes(s,'随机出号只作对照。')
 s=base.body_slide(prs,'资金边界和停止条件','风险控制'); base.card(s,Inches(.78),Inches(1.8),Inches(3.55),Inches(3.9),'每期成本','3个定位胆号码\n每个1元\n每期合计3元','gold',18,22); base.card(s,Inches(4.58),Inches(1.8),Inches(3.55),Inches(3.9),'建议本金',f'主方案{PERIODS*3}元\n随机对照{PERIODS*3}元\n合计{PERIODS*6}元','blue',18,21); base.card(s,Inches(8.38),Inches(1.8),Inches(3.55),Inches(3.9),'停止规则',f'止盈：不设置\n止损：不设置\n每阶段满{PERIODS}期硬停止\n禁止倍投追损','green',18,18); base.fixed.set_notes(s,'金额统一使用元，硬停止替代临时改动。')
 s=base.body_slide(prs,'怎样记录，怎样做结论','验证方法'); base.card(s,Inches(.78),Inches(1.82),Inches(5.55),Inches(3.9),'每期记录','阶段、期号、实际位置、冻结号码、开奖号、是否命中。改码、漏期或并投必须备注。','gold',18,18); base.card(s,Inches(6.58),Inches(1.82),Inches(5.55),Inches(3.9),'运行后再比较','主方案与随机对照使用相同成本和期数，只用真实挂机记录比较。','blue',18,18); base.fixed.add_text(s,Inches(.9),Inches(6.08),Inches(11.5),Inches(.36),'现在只冻结规则，不提前宣布有效或无效。',18,base.COLORS['white'],True,PP_ALIGN.CENTER); base.fixed.set_notes(s,'真正结论来自后续实际记录。')
 base.fixed.add_end(prs,end); base.fixed.save(prs,PPT)


def validate(g):
 errs=[]
 with ZipFile(SCHEME) as z:
  names=z.namelist()
  if sum(n.endswith('.txt') and '01_主方案' in n for n in names)!=3:errs.append('主方案TXT数量错误')
  if sum(n.endswith('.txt') and '02_随机对照' in n for n in names)!=3:errs.append('对照TXT数量错误')
  for n in names:
   if n.endswith('.txt'):
    raw=z.read(n); text=raw.decode('gbk')
    if b'\r\n' not in raw or 'SchemeCreator=\r\n' not in text:errs.append(f'TXT格式错误:{n}')
    if '01_主方案' in n and not text.startswith('True\r\n'):errs.append(f'主方案未启用:{n}')
    if '02_随机对照' in n and not text.startswith('False\r\n'):errs.append(f'对照未关闭:{n}')
 prs=Presentation(PPT); visible='\n'.join(sh.text for s in prs.slides for sh in s.shapes if hasattr(sh,'text') and sh.text)
 if len(prs.slides)!=10:errs.append(f'PPT页数{len(prs.slides)}')
 for i,s in enumerate(prs.slides,1):
  if not s.notes_slide.notes_text_frame.text.strip():errs.append(f'第{i}页无备注')
 for bad in [ID,'SET_001','更严谨的说法','四组测试码预设','为了方便演示']:
  if bad in visible:errs.append(f'PPT禁词:{bad}')
 if re.search(r'\d+(?:\.\d+)?U\b',visible):errs.append('金额单位U')
 for n in ['百位','十位','个位']:
  if ' '.join(map(str,g[n]['冻结号码'])) not in visible:errs.append(f'号码缺失:{n}')
  if [x['数字'] for x in g[n]['后继计数排序'][:3]]!=g[n]['冻结号码']:errs.append(f'排序与冻结号码不一致:{n}')
 if errs:raise ValueError(';'.join(errs))
 return {'PPT页数':10,'隐藏页':0,'演讲者备注':'ALL_SLIDES','TXT':'6_FILES_GBK_CRLF','号码一致性':'PASS','并列排序一致性':'PASS'}


def main():
 if OUT.exists():shutil.rmtree(OUT)
 OUT.mkdir(parents=True); rows=base.parse_draws(INPUT)
 if len(rows)!=200: raise ValueError(f'本批次要求200期，实际{len(rows)}期')
 g=freeze(rows); hist=audit(rows); data_range=f'{rows[0][0]}—{rows[-1][0]}，共{len(rows)}期'
 build_package(g,rows[-1][0],data_range); build_ppt(g,rows[-1][0])
 evidence={'批次ID':'BATCH-B396-FIRST-ORDER-TRANSITION-001','方案内部ID':ID,'自然名称':PROJECT,'阶段':'PRE_RUN_SETUP','数据来源':{'文件':str(INPUT.relative_to(ROOT)).replace('\\','/'),'期号范围':f'{rows[0][0]}—{rows[-1][0]}','总期数':len(rows),'数据复用状态':'REUSED_DATA_NOT_INDEPENDENT_HOLDOUT'},'分析角度':[{'角度ID':'ANG-069','名称':'一阶马尔可夫转移'},{'角度ID':'ANG-072','名称':'条件频率'}],'观察对象':'百位、十位、个位分别建模','计算过程':'同位置一阶转移频次；最新前值为条件；固定并列规则取前3名。','选取规则':'单一预注册模型，无窗口扫描、无参数择优；样本不足5次回退全局频次。','最终结果':g,'软件字段映射':'三份定码轮换TXT；运行前冻结；顶部方案轮投手工开启；运行中不更新。','历史滚动审计':hist,'前向计划':{'主方案期数':PERIODS,'随机对照期数':PERIODS,'每期成本':3,'建议本金':PERIODS*6,'止盈':'不设置','止损':'不设置','替代停止条件':f'每阶段满{PERIODS}期硬停止','倍投':'禁止'},'结论边界':'历史审计只记录旧数据表现，不替代新的实际挂机结果。'}
 EVIDENCE.write_text(json.dumps(evidence,ensure_ascii=False,indent=2)+'\n',encoding='utf-8'); check=validate(g); seo,outer=youtube.build(OUT,SEO)
 files=[SCHEME,PPT,seo,outer,EVIDENCE]; MANIFEST.write_text(json.dumps({'project':PROJECT,'stage':'PRE_RUN_SETUP','internal_id':ID,'validation':check,'files':{p.name:base.sha256(p) for p in files}},ensure_ascii=False,indent=2)+'\n',encoding='utf-8')
 print('B396_DELIVERY_OK',','.join(f"{n}:{g[n]['冻结号码']}" for n in ['百位','十位','个位']),f"walk_forward={hist['总计']['命中次数']}/{hist['总计']['预测次数']}",outer.name)
if __name__=='__main__':main()
