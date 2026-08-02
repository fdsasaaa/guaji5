#!/usr/bin/env python3
from __future__ import annotations

import hashlib
import json
import math
import shutil
import sys
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TOOLS = ROOT / "tools"
sys.path.insert(0, str(TOOLS))
import materialize_ppt_fixed_pages as fixed  # noqa: E402

BATCH_ID = "BATCH-V393-RULE3-VS-RANDOM-001"
SCHEME_ID = "B393-SET-001"
SLUG = "B393_SET_001_规则三码对随机"
INPUT_DATA = ROOT / "01_本次输入" / "哈希分分彩_20260731_0181至0380.txt"
OUT_ROOT = ROOT / "dist" / "B393_SET_001_delivery"
PACKAGE_STAGE = OUT_ROOT / "package"
ZIP_PATH = OUT_ROOT / f"{SLUG}_方案套.zip"
PPT_PATH = OUT_ROOT / f"{SLUG}_人工讲解型PPT.pptx"
MANIFEST_PATH = OUT_ROOT / "DELIVERY_MANIFEST.json"

COLORS = {
    "bg": (8, 12, 17),
    "panel": (18, 25, 32),
    "panel2": (27, 36, 45),
    "white": (244, 247, 250),
    "gray": (166, 178, 188),
    "gold": (232, 177, 64),
    "green": (65, 201, 151),
    "red": (226, 93, 93),
    "blue": (77, 153, 230),
    "line": (68, 82, 96),
}


def sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def parse_draws(path: Path) -> list[tuple[str, list[int]]]:
    rows: list[tuple[str, list[int]]] = []
    for line_no, raw in enumerate(path.read_text(encoding="utf-8").splitlines(), 1):
        line = raw.strip()
        if not line:
            continue
        if "=" not in line:
            raise ValueError(f"开奖数据第{line_no}行缺少等号")
        issue, number = line.split("=", 1)
        if not issue.isdigit() or len(number) != 5 or not number.isdigit():
            raise ValueError(f"开奖数据第{line_no}行格式错误: {line}")
        rows.append((issue, [int(c) for c in number]))
    if len(rows) != 200:
        raise ValueError(f"预期200期，实际{len(rows)}期")
    issues = [issue for issue, _ in rows]
    if issues != sorted(issues):
        raise ValueError("期号不是升序")
    return rows


def simulate(rows: list[tuple[str, list[int]]]) -> dict:
    groups = [(0, 3, 6), (1, 4, 7), (2, 5, 8), (3, 6, 9)]
    group_index = 0
    hits: list[int] = []
    trace: list[dict] = []
    miss_streak = 0
    max_miss_streak = 0
    for issue, digits in rows:
        target = digits[2]
        selected = groups[group_index]
        hit = int(target in selected)
        hits.append(hit)
        trace.append({"期号": issue, "百位": target, "本期组": " ".join(map(str, selected)), "命中": bool(hit)})
        if hit:
            miss_streak = 0
        else:
            miss_streak += 1
            max_miss_streak = max(max_miss_streak, miss_streak)
            group_index = (group_index + 1) % len(groups)

    windows = {"训练": (0, 120), "验证": (120, 160), "审计": (160, 200), "全样本": (0, 200)}
    split_stats = {}
    for name, (start, end) in windows.items():
        count = sum(hits[start:end])
        split_stats[name] = {
            "开始序号": start + 1,
            "结束序号": end,
            "命中": count,
            "总数": end - start,
            "命中率": count / (end - start),
        }

    total_hits = sum(hits)
    p_value = sum(math.comb(200, k) * (0.3 ** k) * (0.7 ** (200 - k)) for k in range(total_hits, 201))
    fixed_controls = {}
    for group in groups:
        result = [int(digits[2] in group) for _, digits in rows]
        fixed_controls[" ".join(map(str, group))] = {
            "命中": sum(result),
            "总数": len(result),
            "命中率": sum(result) / len(result),
        }

    result = {
        "方案ID": SCHEME_ID,
        "批次ID": BATCH_ID,
        "问题": "百位等距三码组按‘连挂1期换号’运行，是否优于固定三码与随机三码？",
        "位置": "百位",
        "分组": [" ".join(map(str, g)) for g in groups],
        "换号规则": "连挂1期换号（GUI标签语义；真实运行仍需软件日志确认）",
        "分段统计": split_stats,
        "全样本命中": total_hits,
        "全样本总数": 200,
        "全样本命中率": total_hits / 200,
        "随机理论命中率": 0.3,
        "单侧二项检验p值": p_value,
        "最大连续未命中": max_miss_streak,
        "固定三码对照": fixed_controls,
        "结论": "未达到统计显著；只能进入软件导入与前向验证，不能宣称存在稳定优势。",
        "执行边界": [
            "三个TXT是主方案与两个对照组，禁止同时并投。",
            "全程平倍，每期每码1U；不得追加马丁。",
            "建议本金300U；单轮20期或亏损60U，先到即停。",
            "止盈不设置；目标是验证规则，不是追求短期盈利。",
        ],
        "追踪样例": trace[:20],
    }
    expected = {"训练": (36, 120), "验证": (11, 40), "审计": (18, 40), "全样本": (65, 200)}
    for name, (hit_count, total) in expected.items():
        got = split_stats[name]
        if (got["命中"], got["总数"]) != (hit_count, total):
            raise AssertionError(f"{name}统计漂移: {got}")
    if max_miss_streak != 11:
        raise AssertionError(f"最大连挂漂移: {max_miss_streak}")
    return result


def common_lines(strategy: str, play_name: str, first_line: str) -> list[str]:
    return [
        first_line, strategy, "软件名称=CXGGJ", "玩法类型=定位胆", f"玩法名称={play_name}", "金额模式=2",
        "投注监控=False-", "投注监控模式=0", "任选中奖=1-10", "任选位置=",
    ]


def tail_lines() -> list[str]:
    return [
        "翻倍方式=0", "正集=True", "倍投类型=0", "倍投计划=1,1,1,1,1,1,1,1,1,1", "倍投方案=",
        "显示更多=False", "真实投注1=False-50000", "真实投注2=False-50000", "模拟投注1=False-50000",
        "模拟投注2=False-50000", "盈利跳转=False-50000-1", "亏损跳转=False-50000-1",
        "盈利停止=False-50000", "亏损停止=False-50000", "投注时间=False", "投注时间类型=0",
        "范围开始时间=False-09:01:00", "范围停止时间=False-21:32:00", "范围停止类型=0",
        "倒计时停止时间=02:00:00", "倒计时停止类型=0",
    ]


def scheme_main() -> str:
    lines = common_lines("定码轮换", "百位", "False") + ["换号规则=3", "换号期数=1"] + tail_lines()
    lines += ["定码轮换内容=0 3 6;1 4 7;2 5 8;3 6 9", "定码轮换单组=False", "SchemeCreator="]
    return "\r\n".join(lines) + "\r\n"


def scheme_fixed() -> str:
    lines = common_lines("定码轮换", "百位", "False") + ["换号规则=9", "换号期数=5"] + tail_lines()
    lines += ["定码轮换内容=0 3 6", "定码轮换单组=True", "SchemeCreator="]
    return "\r\n".join(lines) + "\r\n"


def scheme_random() -> str:
    lines = common_lines("随机出号", "百位", "True") + ["换号规则=0", "换号期数=3"] + tail_lines()
    lines += ["随机出号模板=模板1", "随机出号个数=3", "SchemeCreator="]
    return "\r\n".join(lines) + "\r\n"


def write_gbk(path: Path, text: str) -> None:
    path.write_bytes(text.encode("gbk"))


def build_package(stats: dict) -> None:
    if PACKAGE_STAGE.exists():
        shutil.rmtree(PACKAGE_STAGE)
    PACKAGE_STAGE.mkdir(parents=True)
    write_gbk(PACKAGE_STAGE / "B001_等距三码连挂1期切组-定码轮换.txt", scheme_main())
    write_gbk(PACKAGE_STAGE / "C001_固定036三码对照-定码轮换.txt", scheme_fixed())
    write_gbk(PACKAGE_STAGE / "C002_随机三码对照-随机出号.txt", scheme_random())

    readme = f"""# {SCHEME_ID} 规则三码对随机基准验证套

## 这套方案在验证什么

百位使用四组三码：`0 3 6 → 1 4 7 → 2 5 8 → 3 6 9`。
当前组未命中一次，就切到下一组；命中则继续保留当前组。

## 三个文件怎么用

- `B001`：主方案，等距三码按连挂1期切组。
- `C001`：固定三码对照，全程只投0、3、6。
- `C002`：随机三码对照，由软件每期随机取3码。

**三个文件只能分别运行，禁止同时并投。** 否则成本叠加，对照实验失效。

## 历史结果

- 训练：36/120 = 30.00%
- 验证：11/40 = 27.50%
- 审计：18/40 = 45.00%
- 全部：65/200 = 32.50%
- 随机理论：30.00%
- 单侧二项检验：p = {stats['单侧二项检验p值']:.3f}

32.50%没有达到统计显著，审计段45%也不能单独拿出来宣传。当前结论是：**可测试，但没有证据证明稳定优于随机。**

## 资金与停止

- 全程平倍：每码1U，每期3U。
- 建议本金：300U；若1U=1元，则建议本金300元，按比例换算。
- 单轮：20期。
- 止损：60U，或跑满20期，先到即停。
- 止盈：不设置。
- 禁止马丁、追损、三个文件并投。

## 软件核对

导入后必须检查：玩法=定位胆百位；B001有四段号码；换号方式显示为“连挂1期换号”；三个文件都未加密。若显示与此不一致，立即停止运行并保留截图/日志。
"""
    (PACKAGE_STAGE / "00_使用说明.md").write_text(readme, encoding="utf-8")
    (PACKAGE_STAGE / "01_历史验证.json").write_text(json.dumps(stats, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    funding = {
        "计价单位": "1U=每个号码的单注金额", "每期号码数": 3, "每期暴露": "3U", "建议本金": "300U",
        "单轮期数": 20, "止盈止损模式": "STOP_LOSS_ONLY", "止盈": "不设置", "止损": "60U",
        "达到后动作": "停止本轮，不加码，不当天重启", "总实验上限": "两轮40期，累计毛暴露120U",
    }
    (PACKAGE_STAGE / "02_资金冻结.json").write_text(json.dumps(funding, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    checklist = """# 导入与运行核对表

- [ ] 三个TXT均能正常导入，未显示加密
- [ ] B001玩法为定位胆/百位
- [ ] B001号码分组为036、147、258、369
- [ ] B001换号规则显示为连挂1期换号
- [ ] C001固定为036且从不换号
- [ ] C002每期随机取3码
- [ ] 三个方案分别运行，不并投
- [ ] 倍投始终为1
- [ ] 记录每期选码、开奖结果、命中、当前组和是否切组
- [ ] 20期或亏损60U立即停止
"""
    (PACKAGE_STAGE / "03_导入运行核对表.md").write_text(checklist, encoding="utf-8")
    if ZIP_PATH.exists():
        ZIP_PATH.unlink()
    with ZipFile(ZIP_PATH, "w", ZIP_DEFLATED) as zf:
        for path in sorted(PACKAGE_STAGE.iterdir()):
            zf.write(path, arcname=path.name)


def new_body_slide(prs: Presentation, title: str, kicker: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*COLORS["bg"])
    if kicker:
        fixed.add_text(slide, Inches(0.78), Inches(0.38), Inches(4.0), Inches(0.30), kicker, 12, COLORS["gold"], True)
    fixed.add_text(slide, Inches(0.78), Inches(0.72), Inches(11.7), Inches(0.58), title, 28, COLORS["white"], True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.78), Inches(1.42), Inches(11.74), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*COLORS["line"])
    line.line.fill.background()
    return slide


def add_card(slide, x, y, w, h, title, body, accent="gold", title_size=16, body_size=13):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS["panel"])
    shape.line.color.rgb = RGBColor(*COLORS["line"])
    shape.line.width = Pt(1)
    fixed.add_text(slide, x + Inches(0.24), y + Inches(0.16), w - Inches(0.48), Inches(0.32), title, title_size, COLORS[accent], True)
    fixed.add_text(slide, x + Inches(0.24), y + Inches(0.55), w - Inches(0.48), h - Inches(0.70), body, body_size, COLORS["white"], False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP)
    return shape


def add_metric(slide, x, y, w, label, value, sub="", accent="gold"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, Inches(1.12))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS["panel"])
    shape.line.color.rgb = RGBColor(*COLORS["line"])
    fixed.add_text(slide, x + Inches(0.18), y + Inches(0.12), w - Inches(0.36), Inches(0.22), label, 11, COLORS["gray"])
    fixed.add_text(slide, x + Inches(0.18), y + Inches(0.36), w - Inches(0.36), Inches(0.42), value, 25, COLORS[accent], True)
    if sub:
        fixed.add_text(slide, x + Inches(0.18), y + Inches(0.80), w - Inches(0.36), Inches(0.20), sub, 10, COLORS["gray"])


def add_group_chip(slide, x, y, text, active=False):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(2.35), Inches(0.72))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*(COLORS["gold"] if active else COLORS["panel2"]))
    shape.line.color.rgb = RGBColor(*(COLORS["gold"] if active else COLORS["line"]))
    fixed.add_text(slide, x, y, Inches(2.35), Inches(0.72), text, 22, COLORS["bg"] if active else COLORS["white"], True, PP_ALIGN.CENTER)


def hide_slide(slide) -> None:
    slide._element.set("show", "0")


def build_ppt(stats: dict) -> None:
    cover_png = ROOT / "assets" / "ppt" / "fixed_pages" / "首页背景图谱.png"
    end_png = ROOT / "assets" / "ppt" / "fixed_pages" / "固定最后一页_画面.png"
    for path in [cover_png, end_png]:
        if not path.exists():
            raise FileNotFoundError(f"固定页资源未物化: {path}")

    prs = fixed.new_prs()
    cover = fixed.add_cover(prs, cover_png)
    fixed.add_text(cover, Inches(0.82), Inches(1.05), Inches(8.4), Inches(0.62), "三码轮换能赢随机吗", 32, COLORS["white"], True)
    fixed.add_text(cover, Inches(0.84), Inches(1.82), Inches(7.6), Inches(0.34), "规则组三码 · 固定对照 · 随机对照", 16, COLORS["gold"], True)
    fixed.add_text(cover, Inches(0.84), Inches(5.88), Inches(7.4), Inches(0.34), "200期历史复核｜前向40期验证方案", 14, COLORS["white"])
    fixed.set_notes(cover, "本期只验证一个问题：有规则地轮换三码，是否真的比固定三码和随机三码更好。")
    fixed.build_second_slide(prs)

    slide = new_body_slide(prs, "先把问题说清楚", "研究问题")
    add_card(slide, Inches(0.78), Inches(1.78), Inches(5.55), Inches(3.55), "主方案", "百位每期投3个数字。\n当前组挂1期，就切到下一组；命中则继续用当前组。", "gold", 18, 18)
    add_card(slide, Inches(6.58), Inches(1.78), Inches(5.55), Inches(3.55), "必须同时回答", "它是否优于：\n① 固定投同一组三码\n② 每期随机选三码", "blue", 18, 18)
    fixed.add_text(slide, Inches(0.95), Inches(5.65), Inches(11.0), Inches(0.58), "只看主方案的短期好成绩，没有意义；对照组才决定它是否真的有价值。", 18, COLORS["white"], True, PP_ALIGN.CENTER)
    fixed.set_notes(slide, "不要先问赚了多少。先问这个规则有没有超过同成本的固定和随机基准。")

    slide = new_body_slide(prs, "四组三码，挂一次就切组", "核心规则")
    xs = [0.80, 3.33, 5.86, 8.39]
    for i, text in enumerate(["0 3 6", "1 4 7", "2 5 8", "3 6 9"]):
        add_group_chip(slide, Inches(xs[i]), Inches(2.02), text, active=(i == 0))
        if i < 3:
            fixed.add_text(slide, Inches(xs[i] + 2.28), Inches(2.15), Inches(0.34), Inches(0.36), "→", 22, COLORS["gray"], True, PP_ALIGN.CENTER)
    fixed.add_text(slide, Inches(10.76), Inches(2.15), Inches(0.70), Inches(0.36), "↺", 24, COLORS["gray"], True, PP_ALIGN.CENTER)
    add_card(slide, Inches(0.80), Inches(3.18), Inches(5.25), Inches(2.14), "未命中", "只挂1期，下一期切到右边一组。", "red", 18, 18)
    add_card(slide, Inches(6.30), Inches(3.18), Inches(5.25), Inches(2.14), "命中", "当前组不动，下一期继续使用。", "green", 18, 18)
    fixed.add_text(slide, Inches(0.88), Inches(5.75), Inches(10.95), Inches(0.38), "注意：这里采用软件GUI标签语义；实际切组行为仍要用运行日志确认。", 14, COLORS["gold"], True, PP_ALIGN.CENTER)
    fixed.set_notes(slide, "规则很简单：挂一次就向右切组，命中就停在原组。切组方向固定循环。")

    slide = new_body_slide(prs, "完整走一遍，就不会理解错", "运行案例")
    example = [
        ("第1期", "投 0 3 6", "百位=4", "挂 → 切147", "red"),
        ("第2期", "投 1 4 7", "百位=7", "中 → 留147", "green"),
        ("第3期", "投 1 4 7", "百位=0", "挂 → 切258", "red"),
        ("第4期", "投 2 5 8", "百位=5", "中 → 留258", "green"),
    ]
    y = 1.78
    for period, bet, result, action, accent in example:
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(y), Inches(11.35), Inches(0.86))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*COLORS["panel"])
        card.line.color.rgb = RGBColor(*COLORS["line"])
        fixed.add_text(slide, Inches(1.05), Inches(y + 0.08), Inches(1.25), Inches(0.68), period, 15, COLORS["gold"], True)
        fixed.add_text(slide, Inches(2.35), Inches(y + 0.08), Inches(2.25), Inches(0.68), bet, 18, COLORS["white"], True)
        fixed.add_text(slide, Inches(4.75), Inches(y + 0.08), Inches(2.10), Inches(0.68), result, 16, COLORS["gray"])
        fixed.add_text(slide, Inches(7.05), Inches(y + 0.08), Inches(4.55), Inches(0.68), action, 17, COLORS[accent], True)
        y += 1.03
    fixed.set_notes(slide, "这个例子重点说明两个动作：挂后切组，命中留组。没有其他隐藏判断。")

    slide = new_body_slide(prs, "为什么没有选冷热和重号", "候选审议")
    add_card(slide, Inches(0.78), Inches(1.72), Inches(3.55), Inches(3.85), "冷热温出号", "可导入，但数字类型对应的真实选号含义仍不清楚。\n\n结论：不拿未知语义做正式主线。", "red", 17, 15)
    add_card(slide, Inches(4.50), Inches(1.72), Inches(3.55), Inches(3.85), "重号侦测", "可导入和显示，但最近开出、反集等真实运行行为仍缺日志证据。\n\n结论：暂不作为本批核心。", "red", 17, 15)
    add_card(slide, Inches(8.22), Inches(1.72), Inches(3.55), Inches(3.85), "定码轮换", "分号多段结构已验证；换号规则有明确GUI映射。\n\n结论：选它，因为规则最容易复核。", "green", 17, 15)
    fixed.add_text(slide, Inches(1.15), Inches(5.95), Inches(10.6), Inches(0.38), "能导入不等于能解释；能解释不等于有优势。两道门都要过。", 17, COLORS["white"], True, PP_ALIGN.CENTER)
    fixed.set_notes(slide, "不是功能越多越高级。语义不清的功能会污染实验，所以本批主动克制。")

    slide = new_body_slide(prs, "200期结果：整体只高2.5个百分点", "历史验证")
    add_metric(slide, Inches(0.80), Inches(1.78), Inches(2.72), "训练区", "36 / 120", "30.00%", "gray")
    add_metric(slide, Inches(3.67), Inches(1.78), Inches(2.72), "验证区", "11 / 40", "27.50%", "red")
    add_metric(slide, Inches(6.54), Inches(1.78), Inches(2.72), "审计区", "18 / 40", "45.00%", "green")
    add_metric(slide, Inches(9.41), Inches(1.78), Inches(2.72), "全部", "65 / 200", "32.50%", "gold")
    fixed.add_text(slide, Inches(1.02), Inches(3.45), Inches(2.1), Inches(0.34), "规则组三码", 15, COLORS["white"], True)
    bar1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(3.05), Inches(3.50), Inches(6.50), Inches(0.36))
    bar1.fill.solid(); bar1.fill.fore_color.rgb = RGBColor(*COLORS["gold"]); bar1.line.fill.background()
    fixed.add_text(slide, Inches(9.70), Inches(3.40), Inches(1.20), Inches(0.48), "32.5%", 18, COLORS["gold"], True)
    fixed.add_text(slide, Inches(1.02), Inches(4.35), Inches(2.1), Inches(0.34), "随机理论", 15, COLORS["white"], True)
    bar2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(3.05), Inches(4.40), Inches(6.00), Inches(0.36))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = RGBColor(*COLORS["blue"]); bar2.line.fill.background()
    fixed.add_text(slide, Inches(9.20), Inches(4.30), Inches(1.20), Inches(0.48), "30.0%", 18, COLORS["blue"], True)
    fixed.add_text(slide, Inches(0.95), Inches(5.35), Inches(11.0), Inches(0.70), "审计段45%看起来漂亮，但训练和验证没有同步改善。\n不能把一个高段落当成稳定规律。", 17, COLORS["white"], True, PP_ALIGN.CENTER)
    fixed.set_notes(slide, "全样本是32.5%，只比随机理论高2.5个百分点。分段不稳定，比单个高点更重要。")

    slide = new_body_slide(prs, "统计结论：目前不能说它赢了随机", "随机基准")
    add_metric(slide, Inches(0.90), Inches(1.85), Inches(3.45), "观察命中", "65次", "200期", "gold")
    add_metric(slide, Inches(4.92), Inches(1.85), Inches(3.45), "随机期望", "60次", "30% × 200", "blue")
    add_metric(slide, Inches(8.94), Inches(1.85), Inches(3.45), "单侧p值", f"{stats['单侧二项检验p值']:.3f}", "未达0.05", "red")
    add_card(slide, Inches(1.10), Inches(3.55), Inches(10.65), Inches(1.82), "该怎么解释", "65次命中完全可能由随机波动产生。\n当前状态只能写成‘待前向验证’，不能写成‘验证有效’。", "red", 19, 18)
    fixed.add_text(slide, Inches(1.25), Inches(5.78), Inches(10.35), Inches(0.44), "固定对照、随机对照和主方案必须同成本、同周期、分开运行。", 16, COLORS["green"], True, PP_ALIGN.CENTER)
    fixed.set_notes(slide, "p值0.242，离常用显著门槛0.05很远。结果不是失败，而是诚实地阻止夸大。")

    slide = new_body_slide(prs, "前向怎么跑：一次只开一个方案", "执行与资金")
    add_card(slide, Inches(0.78), Inches(1.72), Inches(3.52), Inches(3.75), "运行顺序", "第一轮：主方案20期\n第二轮：固定036对照20期\n第三轮：随机三码对照20期\n\n每轮独立记录。", "blue", 17, 15)
    add_card(slide, Inches(4.54), Inches(1.72), Inches(3.52), Inches(3.75), "资金冻结", "1U = 每个号码单注金额\n每期3U\n建议本金300U\n止损60U\n止盈不设置", "gold", 17, 15)
    add_card(slide, Inches(8.30), Inches(1.72), Inches(3.52), Inches(3.75), "硬性禁止", "不并投三个方案\n不马丁、不追损\n20期或亏损60U即停\n当天不重启", "red", 17, 15)
    fixed.add_text(slide, Inches(1.00), Inches(5.82), Inches(11.0), Inches(0.46), "300U是实验缓冲，不是收益承诺；赔率未知时，不设置止盈更诚实。", 16, COLORS["white"], True, PP_ALIGN.CENTER)
    fixed.set_notes(slide, "资金以U表示，方便按真实单注金额换算。单轮20期和60U止损是硬边界。")

    slide = new_body_slide(prs, "最终判断：保留测试，不保留幻想", "结论")
    add_card(slide, Inches(0.95), Inches(1.82), Inches(5.25), Inches(3.55), "可以做", "三个TXT都具备真实导入结构。\n规则明确，容易用日志核对。\n适合做一轮前向对照实验。", "green", 19, 17)
    add_card(slide, Inches(6.50), Inches(1.82), Inches(5.25), Inches(3.55), "不能做", "不能宣称历史已证明有效。\n不能只展示审计段45%。\n不能靠倍投把弱选号包装成盈利方案。", "red", 19, 17)
    fixed.add_text(slide, Inches(1.10), Inches(5.68), Inches(10.9), Inches(0.60), "前向40期后，只比较主方案与两个对照的净结果；\n中途不改组、不换规则。", 18, COLORS["gold"], True, PP_ALIGN.CENTER)
    fixed.set_notes(slide, "这套方案的价值是可证伪。前向测试若没有稳定超过对照，就停止跟踪。")

    appendix = new_body_slide(prs, "技术附录｜证据与边界", "隐藏页")
    appendix_text = (
        "数据：202607310181—202607310380，共200期，按期号升序。\n"
        "训练/验证/审计：120/40/40；审计区已被本批使用，不再视为独立样本外。\n"
        "主方案：百位；036→147→258→369；连挂1期切下一组；命中留组。\n"
        "统计：36/120、11/40、18/40；合计65/200=32.50%；随机理论30%；单侧p=0.242。\n"
        "最大连续未命中：11期。\n"
        "软件证据：定码轮换多段分号结构已验证导入；换号规则3映射GUI‘连挂N期换号’。\n"
        "未验证：GUI标签与真实运行状态机是否完全一致；必须以前向日志确认。\n"
        "资金：每期3U；建议本金300U；止盈不设置；止损60U；单轮20期。"
    )
    add_card(appendix, Inches(0.82), Inches(1.65), Inches(11.55), Inches(4.95), "完整记录", appendix_text, "gold", 17, 14)
    fixed.set_notes(appendix, "隐藏附录保存完整数字、证据等级、软件边界和资金冻结依据。")
    hide_slide(appendix)

    fixed.add_end(prs, end_png)
    prs.save(PPT_PATH)
    check = Presentation(PPT_PATH)
    if len(check.slides) != 12:
        raise AssertionError(f"PPT页数错误: {len(check.slides)}")
    if check.slides[10]._element.get("show") != "0":
        raise AssertionError("技术附录未隐藏")
    for idx, slide in enumerate(check.slides, 1):
        if not slide.notes_slide.notes_text_frame.text.strip():
            raise AssertionError(f"第{idx}页缺少演讲者备注")


def build_manifest(stats: dict) -> None:
    files = [ZIP_PATH, PPT_PATH]
    manifest = {
        "scheme_id": SCHEME_ID,
        "batch_id": BATCH_ID,
        "generated_files": {p.name: {"size": p.stat().st_size, "sha256": sha256(p)} for p in files},
        "validation": {
            "data_rows": 200,
            "history_hits": stats["全样本命中"],
            "history_rate": stats["全样本命中率"],
            "p_value": stats["单侧二项检验p值"],
            "ppt_slides": 12,
            "hidden_appendix": True,
            "fixed_cover_second_end": True,
            "scheme_creator_empty": True,
        },
    }
    MANIFEST_PATH.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def main() -> None:
    if OUT_ROOT.exists():
        shutil.rmtree(OUT_ROOT)
    OUT_ROOT.mkdir(parents=True)
    rows = parse_draws(INPUT_DATA)
    stats = simulate(rows)
    build_package(stats)
    build_ppt(stats)
    build_manifest(stats)
    print(
        "DELIVERY_OK",
        f"scheme={SCHEME_ID}",
        f"hits={stats['全样本命中']}/200",
        f"p={stats['单侧二项检验p值']:.3f}",
        f"zip={ZIP_PATH.name}",
        f"ppt={PPT_PATH.name}",
    )


if __name__ == "__main__":
    main()
