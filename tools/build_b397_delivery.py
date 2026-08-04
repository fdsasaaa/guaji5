#!/usr/bin/env python3
from __future__ import annotations
from collections import Counter
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile
import json, math, re, shutil, sys
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

ROOT=Path(__file__).resolve().parents[1]
sys.path.insert(0,str(ROOT/'tools'))
import build_b394_delivery as base

PROJECT='三区冠军三码'; ID='B397-SET-001'; BATCH='BATCH-B397-ZONE-CHAMPION-001'
PERIODS=24; P0=.30; FUNDING_SEQUENCE=[1,1,2,3,2,1,1,1]; UNIT_EXPOSURE=3
INPUT=ROOT/'01_本次输入'/'哈希分分彩_20260731_0181至0380.txt'
OUT=ROOT/'dist'/f'{PROJECT}_交付'; SCHEME_DIR=OUT/f'{PROJECT}_方案文件夹'
PPT=OUT/f'{PROJECT}_挂机前讲解.pptx'; SEO=OUT/f'{PROJECT}_YouTube_SEO.txt'
EVIDENCE=OUT/f'{PROJECT}_验证记录.json'; MANIFEST=OUT/'DELIVERY_MANIFEST.json'; OUTER=OUT/f'{PROJECT}_完整交付.zip'
POS=[(2,'百位'),(3,'十位'),(4,'个位')]; ZONES=[('低区',(0,1,2)),('中区',(3,4,5,6)),('高区',(7,8,9))]
ACC={'百位':'gold','十位':'blue','个位':'green'}

def choose(values,issues=None):
 c=Counter(values); last={d:max((i for i,x in enumerate(values) if x==d),default=-1) for d in range(10)}
 selected=[]; detail=[]
 for name,digits in ZONES:
  ranked=sorted(digits,key=lambda d:(-c[d],-last[d],d)); selected.append(ranked[0])
  detail.append({'分区':name,'范围':list(digits),'排序':[{'数字':d,'出现次数':c[d],'最近位置':last[d],'最近期号':issues[last[d]] if issues and last[d]>=0 else None} for d in ranked],'冠军':ranked[0]})
 return selected,detail

def tail_p(n,hits,p=P0): return min(1.0,sum(math.comb(n,k)*p**k*(1-p)**(n-k) for k in range(hits,n+1)))
def max_miss(values):
 best=run=0
 for hit in values: run=0 if hit else run+1; best=max(best,run)
 return best
def summary(values):
 n=len(values); h=sum(values)
 return {'预测次数':n,'命中次数':h,'命中比例':round(h/n,6),'随机三码理论基准':P0,'相对基准差':round(h/n-P0,6),'独立近似单侧二项P值':round(tail_p(n,h),6),'最大连续未中':max_miss(values)}

def freeze(rows):
 issues=[x for x,_ in rows]; out={}
 for idx,name in POS:
  selected,detail=choose([d[idx] for _,d in rows],issues)
  out[name]={'位置索引':idx,'样本期数':len(rows),'统计起始期号':rows[0][0],'统计截止期号':rows[-1][0],'固定分区':{'低区':[0,1,2],'中区':[3,4,5,6],'高区':[7,8,9]},'分区排序':detail,'冻结号码':selected,'运行中是否更新':False}
 return out

def audit(rows):
 seg={'校准段':[],'验证段':[],'审计段':[]}; per={name:[] for _,name in POS}; all_hits=[]
 for t in range(60,len(rows)):
  for idx,name in POS:
   selected,_=choose([d[idx] for _,d in rows[:t]]); hit=rows[t][1][idx] in selected
   all_hits.append(hit); per[name].append(hit); seg['校准段' if t<120 else '验证段' if t<160 else '审计段'].append(hit)
 return {'方法':'从第61个样本开始逐期扩展窗口；每次只用此前同位置数据，各区取一名频次冠军。','观察构成':{'开始测试期号':rows[60][0],'结束测试期号':rows[-1][0],'测试彩票开奖期数':len(rows)-60,'每期位置观察数':3,'总位置观察数':(len(rows)-60)*3},'总计':summary(all_hits),'分段':{k:summary(v) for k,v in seg.items()},'分位置':{k:summary(v) for k,v in per.items()},'统计边界':'二项P值仅作独立近似参考；位置间和时间上可能相关。','样本边界':'200期数据已被项目复用，不属于新的独立样本外。'}

def set_value(lines,prefix,value):
 for i,line in enumerate(lines):
  if line.startswith(prefix): lines[i]=prefix+value; return
 raise ValueError(f'缺少字段: {prefix}')
def formal_txt(play,digits):
 lines=base.common('定码轮换',play,True)+['换号规则=9',f'换号期数={PERIODS}']+base.tail(); seq=','.join(map(str,FUNDING_SEQUENCE))
 set_value(lines,'倍投类型=','0'); set_value(lines,'倍投计划=',seq); set_value(lines,'倍投方案=',seq)
 lines += [f"定码轮换内容={' '.join(map(str,digits))}",'定码轮换单组=True','SchemeCreator=']
 return '\r\n'.join(lines)+'\r\n'
def build_scheme_folder(groups):
 if SCHEME_DIR.exists(): shutil.rmtree(SCHEME_DIR)
 SCHEME_DIR.mkdir(parents=True)
 for name in ['百位','十位','个位']:
  digits=groups[name]['冻结号码']; visible=''.join(map(str,digits))
  (SCHEME_DIR/f'{name}{visible}-定码轮换.txt').write_bytes(formal_txt(name,digits).encode('gbk'))

def zone_text(item): return '　'.join(f"{r['数字']}：{r['出现次数']}次" for r in sorted(item['排序'],key=lambda x:x['数字']))
def card_counts(slide,groups,pos,y=1.72):
 for i,item in enumerate(groups[pos]['分区排序']):
  body=f"{zone_text(item)}\n\n冠军：{item['冠军']}"
  base.card(slide,Inches(.72+i*4.05),Inches(y),Inches(3.82),Inches(3.80),item['分区'],body,['gold','blue','green'][i],18,17)
 digits=''.join(map(str,groups[pos]['冻结号码']))
 base.fixed.add_text(slide,Inches(.92),Inches(5.85),Inches(11.45),Inches(.46),f'{pos}最终投注：{digits}',22,base.COLORS['white'],True,PP_ALIGN.CENTER)

def build_ppt(groups,first_issue,last_issue,historical):
 cover=ROOT/'assets'/'ppt'/'fixed_pages'/'首页背景图谱.png'; end=ROOT/'assets'/'ppt'/'fixed_pages'/'固定最后一页_画面.png'; prs=base.fixed.new_prs()
 s=base.fixed.add_cover(prs,cover)
 base.fixed.add_text(s,Inches(.82),Inches(1.02),Inches(9.2),Inches(.72),PROJECT,34,base.COLORS['white'],True)
 base.fixed.add_text(s,Inches(.84),Inches(1.86),Inches(9.7),Inches(.38),'低、中、高三区各选一个频次冠军',17,base.COLORS['gold'],True)
 base.fixed.add_text(s,Inches(.84),Inches(5.88),Inches(9.8),Inches(.34),'挂机前规则说明与真实运行验证',14,base.COLORS['white'])
 base.fixed.set_notes(s,'本期先告诉观众投什么，再解释数字来源、软件设置、资金和验证口径。')

 s=base.body_slide(prs,'本轮到底投什么','先看答案')
 for i,pos in enumerate(['百位','十位','个位']):
  digits=''.join(map(str,groups[pos]['冻结号码']))
  base.card(s,Inches(.72+i*4.05),Inches(1.72),Inches(3.82),Inches(3.70),pos,digits,ACC[pos],19,34)
 base.fixed.add_text(s,Inches(.92),Inches(5.72),Inches(11.45),Inches(.52),'每个位置把0—9分成低、中、高三区，再从每一区选出出现次数最多的一个数字。',17,base.COLORS['white'],True,PP_ALIGN.CENTER)
 base.fixed.set_notes(s,'直接读出：百位269、十位037、个位168。观众在第二页就知道本轮具体投什么。')

 s=base.body_slide(prs,'三区冠军怎样产生','计算方法与统计窗口')
 for i,(title,digits,accent) in enumerate([('低区','0 1 2','gold'),('中区','3 4 5 6','blue'),('高区','7 8 9','green')]): base.card(s,Inches(.72+i*4.05),Inches(1.65),Inches(3.82),Inches(1.55),title,digits,accent,17,24)
 base.card(s,Inches(.72),Inches(3.48),Inches(5.70),Inches(2.25),'三个判断层级','① 出现次数更多优先\n② 次数相同，最近出现者优先\n③ 仍然相同，数字较小者优先','gold',17,17)
 base.card(s,Inches(6.66),Inches(3.48),Inches(5.70),Inches(2.25),'本轮真实数据窗口',f'{first_issue}—{last_issue}\n共200期，百、十、个使用同一窗口。\n统计完成后号码冻结，未来24期不再重算。','blue',17,16)
 base.fixed.set_notes(s,'统计窗口来自仓库原始开奖文件，不是推测。正式挂机前一次性冻结，运行24期期间禁止换号。')

 s=base.body_slide(prs,'百位269怎样算出来','完整复算')
 card_counts(s,groups,'百位')
 base.fixed.set_notes(s,'百位低区2以25次胜出；中区6以26次胜出；高区9以21次胜出，因此百位为269。')

 s=base.body_slide(prs,'十位037怎样算出来','完整复算')
 card_counts(s,groups,'十位')
 low=groups['十位']['分区排序'][0]['排序']; by={r['数字']:r for r in low}
 base.fixed.add_text(s,Inches(.92),Inches(6.36),Inches(11.45),Inches(.30),f"低区0与1同为22次；0最近出现于{by[0]['最近期号']}，比1的{by[1]['最近期号']}更近，所以选0。",12,base.COLORS['gold'],True,PP_ALIGN.CENTER)
 base.fixed.set_notes(s,'十位低区出现并列：0和1都是22次。按第二优先级比较最近出现，0胜出；中区3胜出，高区7胜出。')

 s=base.body_slide(prs,'个位168怎样算出来','完整复算')
 card_counts(s,groups,'个位')
 base.fixed.set_notes(s,'个位低区1为17次；中区6为29次；高区8为25次，所以个位为168。')

 s=base.body_slide(prs,'三份TXT怎样轮流运行','软件设置')
 for i,pos in enumerate(['百位','十位','个位']):
  digits=''.join(map(str,groups[pos]['冻结号码']))
  base.card(s,Inches(.72+i*4.05),Inches(1.72),Inches(3.82),Inches(2.15),f'第{i+1}份｜{pos}',f'{pos}{digits}-定码轮换.txt\n投注：{digits}',ACC[pos],17,18)
  if i<2: base.fixed.add_text(s,Inches(4.42+i*4.05),Inches(2.48),Inches(.42),Inches(.45),'→',24,base.COLORS['gray'],True,PP_ALIGN.CENTER)
 base.card(s,Inches(.72),Inches(4.20),Inches(11.85),Inches(1.58),'必须开启“方案轮投”','轮投开启：每期开奖只运行方案列表中的1份TXT，不是三个位置同时投注。\n先后顺序取决于软件导入后的方案列表；首次导入后要核对前三期实际位置。','gold',17,16)
 base.fixed.add_text(s,Inches(.95),Inches(6.10),Inches(11.35),Inches(.36),'当前目标列表顺序：百位269 → 十位037 → 个位168 → 再回百位；若软件排序不同，以列表实际顺序记录。',13,base.COLORS['red'],True,PP_ALIGN.CENTER)
 base.fixed.set_notes(s,'已核实顶部轮投的核心含义是一期开奖只运行一份方案。TXT本身不锁定列表先后，因此首次导入必须核对软件显示顺序。')

 s=base.body_slide(prs,'每期投多少钱，108元怎样算','资金序列与停止')
 base.card(s,Inches(.72),Inches(1.68),Inches(3.72),Inches(2.08),'1倍基础成本','每期只运行1个位置\n该位置投3个数字\n每个数字1元\n所以1倍＝3元','gold',17,17)
 base.card(s,Inches(4.80),Inches(1.68),Inches(7.57),Inches(2.08),'8步资金序列','1 → 1 → 2 → 3 → 2 → 1 → 1 → 1\n对应：3 → 3 → 6 → 9 → 6 → 3 → 3 → 3元','blue',17,20)
 base.card(s,Inches(.72),Inches(4.08),Inches(5.72),Inches(1.70),'毛投注预算','8步倍数合计12：3×12＝36元\n24期运行3轮：36×3＝108元','green',17,18)
 base.card(s,Inches(6.68),Inches(4.08),Inches(5.69),Inches(1.70),'为什么不是机械追损','最高只到3倍，随后主动降回1倍；8步结束重新开始，三份方案轮投合计24期开奖后强制停止。','red',17,16)
 base.fixed.add_text(s,Inches(.95),Inches(6.14),Inches(11.35),Inches(.36),'108元是完整24期每一步都执行时的毛投注总额，不是预计亏损，也不是盈利目标；资金序列不会提高号码命中概率。',13,base.COLORS['white'],True,PP_ALIGN.CENTER)
 base.fixed.set_notes(s,'逐步算出3元、36元和108元。若没有开启轮投而是三个位置同时投注，1倍将变成9元，当前108元预算立即失效。禁止临时加倍、延长追投、超过3倍或中途换号。')

 s=base.body_slide(prs,'怎样判断方案有效或无效','历史结果与新24期验证')
 total=historical['总计']; obs=historical['观察构成']
 base.card(s,Inches(.72),Inches(1.62),Inches(5.72),Inches(2.15),'420次观察从哪里来',f"从第61个样本开始滚动测试，共{obs['测试彩票开奖期数']}期开奖。\n每期开奖分别观察百、十、个3个位置：\n{obs['测试彩票开奖期数']} × 3 ＝ {obs['总位置观察数']}次位置观察。",'gold',17,17)
 base.card(s,Inches(6.68),Inches(1.62),Inches(5.69),Inches(2.15),'旧数据表现',f"命中：{total['命中次数']} / {total['预测次数']}\n命中率：{total['命中比例']*100:.2f}%\n随机三码理论基准：30%\n只高约{total['相对基准差']*100:.2f}个百分点，未证明优势。",'blue',17,18)
 base.card(s,Inches(.72),Inches(4.08),Inches(11.65),Inches(1.62),'真正的判断标准','冻结规则后记录新的24期开奖：期号、运行位置、数字、倍数、开奖号、命中、投入、返奖、累计盈亏、最大连挂和最大资金暴露。最终分为优于随机／接近随机／低于随机；不能只凭盈利或亏损判断号码规则。','green',17,15)
 base.fixed.add_text(s,Inches(.95),Inches(6.02),Inches(11.35),Inches(.42),'旧数据接近随机基准，目前只能作为候选实验规则，不能作为盈利结论。',16,base.COLORS['red'],True,PP_ALIGN.CENTER)
 base.fixed.set_notes(s,'130/420是140期开奖乘以三个位置。30.95%和30%很接近，不能包装为有效。真正结论只来自冻结后的新24期。')

 s=base.fixed.add_end(prs,end)
 base.card(s,Inches(5.72),Inches(.68),Inches(6.82),Inches(5.80),'挂机前最终确认',f"✓ 统计窗口：{first_issue}—{last_issue}，共200期\n✓ 号码冻结：百位269｜十位037｜个位168\n✓ 开启方案轮投，每期开奖只运行1份TXT\n✓ 首次导入核对软件列表顺序与倍数步进\n✓ 1倍3元｜8步36元｜24期毛投注108元\n✓ 最高3倍；三份方案轮投合计24期开奖后停止\n✓ 不换号、不临时加倍、不延长追投\n✓ 结束后再与30%随机基准比较\n\n本次只冻结规则，不预设结论，24期以后用真实记录说话。",'gold',20,14)
 base.fixed.set_notes(s,'结束前逐项确认：轮投已开启、列表顺序和倍数步进已核对、24期停止。通过网址或Telegram继续了解。')
 base.fixed.save(prs,PPT)

def build_seo(groups):
 nums=f"百位{''.join(map(str,groups['百位']['冻结号码']))}、十位{''.join(map(str,groups['十位']['冻结号码']))}、个位{''.join(map(str,groups['个位']['冻结号码']))}"
 SEO.write_text(f'''标题：三区冠军三码验证：低中高各取1码，24期实测\n\n标签：彩票实验室,三区冠军,定位胆,时时彩研究,彩票数据分析,挂机方案,倍投策略,号码验证,彩票方案测试,五位数彩票\n\n描述：本期验证“三区冠军三码”方案。把0—9固定分为低区、中区和高区，每个位置各选一名历史频次冠军，本轮投注数字为{nums}。资金路径采用1,1,2,3,2,1,1,1，最高3倍，三份方案轮投合计运行24期开奖后停止。内容仅用于数据实验与软件方案验证，不承诺盈利，也不把历史频次解释为未来必然规律。\n''',encoding='utf-8')

def validate(groups):
 errors=[]; seq=','.join(map(str,FUNDING_SEQUENCE)); expected=[]
 for pos in ['百位','十位','个位']:
  digits=groups[pos]['冻结号码']; visible=''.join(map(str,digits)); path=SCHEME_DIR/f'{pos}{visible}-定码轮换.txt'; expected.append(path)
  if not path.exists(): errors.append(f'缺少方案文件: {path.name}'); continue
  raw=path.read_bytes(); text=raw.decode('gbk')
  if b'\r\n' not in raw or not text.startswith('True\r\n'): errors.append(f'TXT编码、换行或启用状态错误: {path.name}')
  if f"定码轮换内容={' '.join(map(str,digits))}\r\n" not in text: errors.append(f'投注数字缺失: {path.name}')
  if f'倍投计划={seq}\r\n' not in text or f'倍投方案={seq}\r\n' not in text: errors.append(f'资金序列不一致: {path.name}')
  if re.search(r'倍投计划=1(?:,1)+\r\n',text): errors.append(f'仍为机械平倍: {path.name}')
 if sorted(p for p in SCHEME_DIR.rglob('*') if p.is_file())!=sorted(expected): errors.append('方案文件夹必须只包含3份可导入TXT')
 prs=Presentation(PPT); visible='\n'.join(sh.text for sl in prs.slides for sh in sl.shapes if hasattr(sh,'text') and sh.text)
 if len(prs.slides)!=10: errors.append(f'PPT页数错误: {len(prs.slides)}')
 for i,sl in enumerate(prs.slides,1):
  if not sl.notes_slide.notes_text_frame.text.strip(): errors.append(f'第{i}页无备注')
 required=['269','037','168','202607310181','202607310380','共200期','0与1同为22次','每期开奖只运行','1 → 1 → 2 → 3 → 2 → 1 → 1 → 1','3×12＝36元','36×3＝108元','130 / 420','30.95%','随机三码理论基准：30%','24期以后用真实记录说话']
 for x in required:
  if x not in visible: errors.append(f'PPT缺少关键内容: {x}')
 for x in [ID,'SET_001','E2','固定参数边界','运行语义','暴露节奏','保证盈利','稳定盈利','回本']:
  if x in visible: errors.append(f'PPT禁词: {x}')
 seo=SEO.read_text(encoding='utf-8')
 if seo.count('标题：')!=1 or seo.count('标签：')!=1 or seo.count('描述：')!=1: errors.append('SEO文件结构错误')
 tags=[x.strip() for x in next((l for l in seo.splitlines() if l.startswith('标签：')),'').removeprefix('标签：').split(',') if x.strip()]
 if not 8<=len(tags)<=10: errors.append(f'SEO标签数量错误: {len(tags)}')
 if OUTER.exists(): OUTER.unlink()
 with ZipFile(OUTER,'w',ZIP_DEFLATED) as z:
  for p in sorted(SCHEME_DIR.rglob('*')):
   if p.is_file(): z.write(p,arcname=f'{SCHEME_DIR.name}/{p.relative_to(SCHEME_DIR).as_posix()}')
  z.write(PPT,arcname=PPT.name); z.write(SEO,arcname=SEO.name)
 with ZipFile(OUTER) as z:
  names=z.namelist(); roots={n.split('/',1)[0] for n in names}
  if roots!={SCHEME_DIR.name,PPT.name,SEO.name}: errors.append(f'完整包根目录不是严格3项: {sorted(roots)}')
  if any(n.lower().endswith(('.md','.csv','.json','.zip')) for n in names): errors.append('完整包含说明、记录、JSON或嵌套ZIP')
 if errors: raise ValueError(';'.join(errors))
 return {'外层ZIP根目录':'EXACTLY_3_ITEMS','方案文件夹':'3_IMPORTABLE_TXT_ONLY','投注数字':'FILENAME_TXT_PPT_VISIBLE','统计窗口':'200_PERIODS_202607310181_TO_202607310380','完整复算':'ALL_THREE_POSITIONS_ALL_DIGITS_VISIBLE','轮投说明':'ONE_TXT_PER_DRAW_LIST_ORDER_CHECK_REQUIRED','资金路径':FUNDING_SEQUENCE,'最高倍数':3,'PPT页数':10,'SEO':'ONE_TITLE_ONE_TAG_LINE_ONE_DESCRIPTION'}

def main():
 if OUT.exists(): shutil.rmtree(OUT)
 OUT.mkdir(parents=True); rows=base.parse_draws(INPUT)
 if len(rows)!=200: raise ValueError('本批要求200期')
 groups=freeze(rows); historical=audit(rows); build_scheme_folder(groups); build_ppt(groups,rows[0][0],rows[-1][0],historical); build_seo(groups); checks=validate(groups)
 stage=UNIT_EXPOSURE*sum(FUNDING_SEQUENCE)*(PERIODS//len(FUNDING_SEQUENCE))
 evidence={'批次ID':BATCH,'方案内部ID':ID,'自然名称':PROJECT,'阶段':'PRE_RUN_SETUP','数据来源':{'文件':str(INPUT.relative_to(ROOT)).replace('\\','/'),'期号范围':f'{rows[0][0]}—{rows[-1][0]}','总期数':len(rows),'三个位置是否同一窗口':True,'数据复用状态':'REUSED_DATA_NOT_INDEPENDENT_HOLDOUT'},'观察对象':'百位、十位、个位分别建模','分析角度':'固定分区频次冠军','计算过程':'0—2、3—6、7—9三区固定；区内频次降序，同频按最近出现与数字升序；每区取1码。','冻结结果':groups,'历史滚动审计':historical,'正式执行':{'顶部方案轮投':'人工勾选','轮投含义':'每期开奖只运行方案列表中的一份TXT','先后顺序来源':'软件导入后的方案列表顺序','首次导入动作':'核对前三次运行位置和倍数步进','运行期数口径':'三份方案轮投合计24期开奖','运行期数':PERIODS,'每期基础成本':UNIT_EXPOSURE,'资金路径类型':'CONTROLLED_PRESSURE_RELEASE','资金序列':FUNDING_SEQUENCE,'最高倍数':3,'8步总倍数':sum(FUNDING_SEQUENCE),'24期毛暴露预算':stage,'建议本金':stage,'止盈':'不设置','止损':'不设置','替代停止条件':'满24期硬停止','软件证据等级':'E2','首次导入核对':True},'结论边界':'资金路径只改变投入节奏，不提高号码中奖概率；旧数据30.95%接近30%随机基准，最终结论只使用新的24期实际挂机记录。'}
 EVIDENCE.write_text(json.dumps(evidence,ensure_ascii=False,indent=2)+'\n',encoding='utf-8')
 files=[PPT,SEO,EVIDENCE,OUTER,*sorted(SCHEME_DIR.glob('*.txt'))]
 MANIFEST.write_text(json.dumps({'project':PROJECT,'stage':'PRE_RUN_SETUP','validation':checks,'external_delivery':OUTER.name,'files':{str(p.relative_to(OUT)):base.sha256(p) for p in files}},ensure_ascii=False,indent=2)+'\n',encoding='utf-8')
 print('B397_DELIVERY_V3_OK',','.join(f"{n}:{groups[n]['冻结号码']}" for n in ['百位','十位','个位']),f'window={rows[0][0]}-{rows[-1][0]}',f'funding={FUNDING_SEQUENCE}',OUTER.name)
if __name__=='__main__': main()
