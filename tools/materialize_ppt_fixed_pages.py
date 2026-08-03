#!/usr/bin/env python3
from __future__ import annotations
from pathlib import Path
import base64, hashlib, json
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches, Pt

ROOT=Path(__file__).resolve().parents[1]
ASSET_DIR=ROOT/'assets'/'ppt'/'fixed_pages'
SOURCE_DIR=ASSET_DIR/'source'
ASSET_DIR.mkdir(parents=True,exist_ok=True)
W,H=Inches(13.333333),Inches(7.5)

def read_b64_parts(stem:str)->bytes:
    parts=sorted(SOURCE_DIR.glob(f'{stem}.*.b64'))
    if not parts: raise FileNotFoundError(f'missing source chunks: {stem}.*.b64')
    return base64.b64decode(''.join(p.read_text(encoding='ascii').strip() for p in parts))

def sha256(path:Path)->str:
    h=hashlib.sha256()
    with path.open('rb') as f:
        for chunk in iter(lambda:f.read(1024*1024),b''): h.update(chunk)
    return h.hexdigest()

def new_prs()->Presentation:
    prs=Presentation(); prs.slide_width=W; prs.slide_height=H
    while len(prs.slides):
        rid=prs.slides._sldIdLst[0].rId; prs.part.drop_rel(rid); del prs.slides._sldIdLst[0]
    return prs

def set_picture_background(slide,image_path:Path)->None:
    pic=slide.shapes.add_picture(str(image_path),0,0,width=W,height=H)
    rid=pic._element.blipFill.blip.rEmbed
    pic._element.getparent().remove(pic._element)
    bg=parse_xml(f'<p:bg {nsdecls("p","a","r")}><p:bgPr><a:blipFill dpi="0" rotWithShape="1"><a:blip r:embed="{rid}"/><a:stretch><a:fillRect/></a:stretch></a:blipFill><a:effectLst/></p:bgPr></p:bg>')
    slide._element.insert(0,bg)

def add_text(slide,x,y,w,h,text,size,color,bold=False,align=PP_ALIGN.LEFT,valign=MSO_ANCHOR.MIDDLE):
    box=slide.shapes.add_textbox(x,y,w,h)
    tf=box.text_frame; tf.clear(); tf.word_wrap=True; tf.vertical_anchor=valign
    para=tf.paragraphs[0]; para.alignment=align
    run=para.add_run(); run.text=text; run.font.name="Microsoft YaHei"; run.font.size=Pt(size); run.font.bold=bold; run.font.color.rgb=RGBColor(*color)
    return box

def set_notes(slide,text):
    slide.notes_slide.notes_text_frame.text=text

def add_link_overlay(slide,x,y,w,h,url):
    s=slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,x,y,w,h)
    s.fill.background(); s.line.fill.background(); s.click_action.hyperlink.address=url

def add_cover(prs,cover_png):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); set_picture_background(slide,cover_png)
    slide.notes_slide.notes_text_frame.text='用自然中文介绍本期待验证主题；不得出现内部编号，也不提前讲结果。'
    return slide

def add_end(prs,end_png):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); set_picture_background(slide,end_png)
    add_link_overlay(slide,Inches(.72),Inches(3.38),Inches(4.55),Inches(.76),'http://www.laocaimi.org')
    add_link_overlay(slide,Inches(.72),Inches(4.26),Inches(4.55),Inches(.62),'https://t.me/laocaimi1314')
    slide.notes_slide.notes_text_frame.text='简短收尾，提醒观众通过网址或Telegram继续了解。'
    return slide

def save(prs,path): path.parent.mkdir(parents=True,exist_ok=True); prs.save(path)

def main():
    cover_jpg=ASSET_DIR/'首页背景图谱_嵌入源.jpg'; end_jpg=ASSET_DIR/'固定最后一页_嵌入源.jpg'
    cover_jpg.write_bytes(read_b64_parts('cover')); end_jpg.write_bytes(read_b64_parts('end'))
    cover_png=ASSET_DIR/'首页背景图谱.png'; end_png=ASSET_DIR/'固定最后一页_画面.png'
    Image.open(cover_jpg).convert('RGB').save(cover_png,optimize=True)
    Image.open(end_jpg).convert('RGB').save(end_png,optimize=True)
    cover=new_prs(); add_cover(cover,cover_png); save(cover,ASSET_DIR/'固定首页模板.pptx')
    end=new_prs(); add_end(end,end_png); save(end,ASSET_DIR/'固定最后一页_标准版.pptx')
    preview=new_prs(); add_cover(preview,cover_png); add_end(preview,end_png); save(preview,ASSET_DIR/'PPT固定首页末页模板_V3.9.4.pptx')
    files=[cover_jpg,end_jpg,cover_png,end_png,ASSET_DIR/'固定首页模板.pptx',ASSET_DIR/'固定最后一页_标准版.pptx',ASSET_DIR/'PPT固定首页末页模板_V3.9.4.pptx']
    manifest={'version':'V3.9.4-FIXED-COVER-END-1','materializer':'tools/materialize_ppt_fixed_pages.py','fixed_second_page':False,'files':{str(p.relative_to(ROOT)).replace('\\','/'):sha256(p) for p in files}}
    (ASSET_DIR/'PPT固定页资源清单.json').write_text(json.dumps(manifest,ensure_ascii=False,indent=2)+'\n',encoding='utf-8')
    print('PPT_FIXED_PAGES_MATERIALIZED',len(files),'second=DISABLED')
if __name__=='__main__': main()
