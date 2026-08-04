#!/usr/bin/env python3
from __future__ import annotations

from collections import Counter
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile
import hashlib
import json
import re
import shutil
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "tools"))
import materialize_ppt_fixed_pages as fixed

INTERNAL_SCHEME_ID = "B395-SET-001"
INPUT_DATA = ROOT / "01_本次输入" / "哈希分分彩_20260731_0181至0380.txt"
OUT_ROOT = ROOT / "dist" / "邻位热差迁移_交付"
PACKAGE = OUT_ROOT / "package"
ZIP_PATH = OUT_ROOT / "邻位热差迁移_方案套.zip"
PPT_PATH = OUT_ROOT / "邻位热差迁移_挂机前讲解.pptx"
MANIFEST = OUT_ROOT / "DELIVERY_MANIFEST.json"

LOOKBACK = 20
RUN_PERIODS = 24
ROUTES = [
    ("万到千", "千位", 0, 1),
    ("千到百", "百位", 1, 2),
    ("百到十", "十位", 2, 3),
    ("十到个", "个位", 3, 4),
]
COLORS = {
    "bg": (8, 12, 17),
    "panel": (18, 25, 32),
    "white": (244, 247, 250),
    "gray": (166, 178, 188),
    "gold": (232, 177, 64),
    "green": (65, 201, 151),
    "red": (226, 93, 93),
    "blue": (77, 153, 230),
    "line": (68, 82, 96),
}


def sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def parse_draws(path: Path):
    rows = []
    for raw in path.read_text(encoding="utf-8").splitlines():
        if not raw.strip():
            continue
        issue, number = raw.strip().split("=", 1)
        if not issue.isdigit() or len(number) != 5 or not number.isdigit():
            raise ValueError(f"开奖格式错误: {raw}")
        rows.append((issue, [int(x) for x in number]))
    if len(rows) < LOOKBACK:
        raise ValueError(f"开奖数据不足{LOOKBACK}期")
    return rows


def gap_group(history, source_pos: int, target_pos: int):
    source = [digits[source_pos] for _, digits in history]
    target = [digits[target_pos] for _, digits in history]
    source_count = Counter(source)
    target_count = Counter(target)
    recency = {
        digit: max((i for i, value in enumerate(source) if value == digit), default=-1)
        for digit in range(10)
    }
    ranking = sorted(
        range(10),
        key=lambda digit: (
            -(source_count[digit] - target_count[digit]),
            -source_count[digit],
            target_count[digit],
            -recency[digit],
            digit,
        ),
    )
    chosen = tuple(sorted(ranking[:3]))
    evidence = [
        {
            "digit": digit,
            "source_count": source_count[digit],
            "target_count": target_count[digit],
            "gap": source_count[digit] - target_count[digit],
        }
        for digit in ranking
    ]
    return chosen, evidence


def freeze_groups(rows):
    history = rows[-LOOKBACK:]
    groups = []
    evidence = {}
    for route, play, source_pos, target_pos in ROUTES:
        chosen, detail = gap_group(history, source_pos, target_pos)
        groups.append(chosen)
        evidence[route] = {
            "play": play,
            "digits": list(chosen),
            "ranking": detail,
        }
    return groups, evidence


def rolling_audit(rows):
    phase_results = {}
    for phase in range(4):
        hits = 0
        trials = 0
        cursor = LOOKBACK
        phase_blocks = []
        while cursor + RUN_PERIODS <= len(rows):
            history = rows[cursor - LOOKBACK:cursor]
            groups = [
                gap_group(history, source_pos, target_pos)[0]
                for _, _, source_pos, target_pos in ROUTES
            ]
            block_hits = 0
            for offset in range(RUN_PERIODS):
                route_index = (offset + phase) % 4
                target_pos = ROUTES[route_index][3]
                actual = rows[cursor + offset][1][target_pos]
                if actual in groups[route_index]:
                    block_hits += 1
            hits += block_hits
            trials += RUN_PERIODS
            phase_blocks.append(block_hits)
            cursor += RUN_PERIODS
        phase_results[str(phase + 1)] = {
            "hits": hits,
            "trials": trials,
            "rate": round(hits / trials, 6) if trials else None,
            "block_hits": phase_blocks,
        }
    return {
        "method": "20期冻结后24期轮投；四种起点全部保留",
        "data_status": "REUSED_DATA_INTERNAL_SANITY_ONLY",
        "random_theory": 0.30,
        "phase_results": phase_results,
        "not_for_public_result_claim": True,
    }


def common(strategy: str, play: str, enabled: bool):
    return [
        "True" if enabled else "False",
        strategy,
        "软件名称=CXGGJ",
        "玩法类型=定位胆",
        f"玩法名称={play}",
        "金额模式=2",
        "投注监控=False-",
        "投注监控模式=0",
        "任选中奖=1-10",
        "任选位置=",
    ]


def tail():
    return [
        "翻倍方式=0",
        "正集=True",
        "倍投类型=0",
        "倍投计划=1,1,1,1,1,1,1,1,1,1",
        "倍投方案=1,1,1,1,1,1,1,1,1,1",
        "显示更多=False",
        "真实投注1=False-50000",
        "真实投注2=False-50000",
        "模拟投注1=False-50000",
        "模拟投注2=False-50000",
        "盈利跳转=False-50000-1",
        "亏损跳转=False-50000-1",
        "盈利停止=False-50000",
        "亏损停止=False-50000",
        "投注时间=False",
        "投注时间类型=0",
        "范围开始时间=False-09:01:00",
        "范围停止时间=False-21:32:00",
        "范围停止类型=0",
        "倒计时停止时间=02:00:00",
        "倒计时停止类型=0",
    ]


def fixed_txt(play: str, digits, enabled: bool = True):
    lines = (
        common("定码轮换", play, enabled)
        + ["换号规则=9", f"换号期数={RUN_PERIODS}"]
        + tail()
        + [
            f"定码轮换内容={' '.join(map(str, digits))}",
            "定码轮换单组=True",
            "SchemeCreator=",
        ]
    )
    return "\r\n".join(lines) + "\r\n"


def random_txt(play: str):
    lines = (
        common("随机出号", play, False)
        + ["换号规则=10", "换号期数=1"]
        + tail()
        + [
            "随机出号模板=模板1",
            "随机出号个数=3",
            "SchemeCreator=",
        ]
    )
    return "\r\n".join(lines) + "\r\n"


def build_package(groups, issue: str):
    if PACKAGE.exists():
        shutil.rmtree(PACKAGE)
    PACKAGE.mkdir(parents=True)

    main_files = [
        ("左热右冷-万到千.txt", "千位", groups[0]),
        ("左热右冷-千到百.txt", "百位", groups[1]),
        ("左热右冷-百到十.txt", "十位", groups[2]),
        ("左热右冷-十到个.txt", "个位", groups[3]),
    ]
    for name, play, digits in main_files:
        (PACKAGE / name).write_bytes(fixed_txt(play, digits).encode("gbk"))

    control_files = [
        ("随机对照-千位.txt", "千位"),
        ("随机对照-百位.txt", "百位"),
        ("随机对照-十位.txt", "十位"),
        ("随机对照-个位.txt", "个位"),
    ]
    for name, play in control_files:
        (PACKAGE / name).write_bytes(random_txt(play).encode("gbk"))

    group_lines = "\n".join(
        f"- {route}，投{play}：{' '.join(map(str, groups[i]))}"
        for i, (route, play, _, _) in enumerate(ROUTES)
    )
    readme = f"""# 邻位热差迁移｜挂机前说明

本方案处于“挂机前验证准备”阶段。它只冻结规则、号码、运行方式和判定口径，不提前判断盈利、亏损或是否优于随机。

## 号码怎样产生

读取截止到 {issue} 的最近{LOOKBACK}期数据。对每一条相邻位置路线，比较同一个数字在左侧来源位置和右侧目标位置的出现次数，优先选择“左侧出现更频繁、右侧出现更少”的3个数字。

{group_lines}

## 运行方式

导入四份主方案后，手工开启软件顶部“方案轮投”。每期只运行一份方案，连续观察{RUN_PERIODS}期。第一期实际从哪一份开始，以软件显示为准并记录。运行中不重启、不改码、不调整顺序。

四份随机对照必须另行运行{RUN_PERIODS}期，不能与主组同时并投。

## 资金边界

- 每期只运行一个位置、3个号码，每个号码按1元计算，单期3元。
- 主组{RUN_PERIODS}期毛投入{RUN_PERIODS * 3}元。
- 随机对照{RUN_PERIODS}期毛投入{RUN_PERIODS * 3}元。
- 建议准备{RUN_PERIODS * 6}元。
- 止盈：不设置。
- 止损：不设置。
- 替代停止条件：主组和对照各运行{RUN_PERIODS}期后立即停止，不追损、不临时换号。

## 必须记录

记录期号、实际方案名称、目标位置、三枚号码、开奖号、是否命中、软件是否重启。完成真实挂机后，再根据实际记录制作结果复盘。
"""
    (PACKAGE / "00_使用说明.md").write_text(readme, encoding="utf-8")

    checklist = f"""# 导入与运行核对表

- [ ] 四份主方案可导入并默认勾选
- [ ] 四份随机对照可导入但默认不勾选
- [ ] 顶部“方案轮投”已手工开启
- [ ] 第一期实际起点已记录
- [ ] 每期只运行一份方案
- [ ] 每个号码按1元计算，全程平倍
- [ ] 主组连续运行{RUN_PERIODS}期后停止
- [ ] 随机对照另行运行{RUN_PERIODS}期
- [ ] 运行中未改码、未重排、未重启
- [ ] 每期实际方案、开奖号和命中情况已记录
"""
    (PACKAGE / "01_导入运行核对表.md").write_text(checklist, encoding="utf-8")

    record = "期号,实际方案名称,来源位置,目标位置,号码,开奖号,是否命中,是否重启,备注\n"
    (PACKAGE / "02_实际挂机记录模板.csv").write_text(record, encoding="utf-8-sig")

    if ZIP_PATH.exists():
        ZIP_PATH.unlink()
    with ZipFile(ZIP_PATH, "w", ZIP_DEFLATED) as zf:
        for path in sorted(PACKAGE.iterdir()):
            zf.write(path, arcname=path.name)


def body_slide(prs, title: str, kicker: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*COLORS["bg"])
    if kicker:
        fixed.add_text(
            slide, Inches(0.78), Inches(0.38), Inches(4), Inches(0.3),
            kicker, 12, COLORS["gold"], True
        )
    fixed.add_text(
        slide, Inches(0.78), Inches(0.72), Inches(11.7), Inches(0.58),
        title, 28, COLORS["white"], True
    )
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.78), Inches(1.42), Inches(11.74), Inches(0.025)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*COLORS["line"])
    line.line.fill.background()
    return slide


def card(slide, x, y, w, h, title, body, accent="gold", title_size=17, body_size=16):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS["panel"])
    shape.line.color.rgb = RGBColor(*COLORS["line"])
    fixed.add_text(
        slide, x + Inches(0.24), y + Inches(0.16), w - Inches(0.48), Inches(0.36),
        title, title_size, COLORS[accent], True
    )
    fixed.add_text(
        slide, x + Inches(0.24), y + Inches(0.58), w - Inches(0.48), h - Inches(0.76),
        body, body_size, COLORS["white"], False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP
    )


def metric(slide, x, y, w, label, value, sub="", accent="gold"):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, Inches(1.15)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS["panel"])
    shape.line.color.rgb = RGBColor(*COLORS["line"])
    fixed.add_text(
        slide, x + Inches(0.18), y + Inches(0.12), w - Inches(0.36), Inches(0.22),
        label, 11, COLORS["gray"]
    )
    fixed.add_text(
        slide, x + Inches(0.18), y + Inches(0.36), w - Inches(0.36), Inches(0.42),
        value, 24, COLORS[accent], True
    )
    if sub:
        fixed.add_text(
            slide, x + Inches(0.18), y + Inches(0.82), w - Inches(0.36), Inches(0.2),
            sub, 10, COLORS["gray"]
        )


def build_ppt(groups, issue: str):
    cover_png = ROOT / "assets" / "ppt" / "fixed_pages" / "首页背景图谱.png"
    end_png = ROOT / "assets" / "ppt" / "fixed_pages" / "固定最后一页_画面.png"
    prs = fixed.new_prs()

    cover = fixed.add_cover(prs, cover_png)
    fixed.add_text(
        cover, Inches(0.82), Inches(1.05), Inches(8.8), Inches(0.62),
        "邻位热差迁移", 32, COLORS["white"], True
    )
    fixed.add_text(
        cover, Inches(0.84), Inches(1.82), Inches(8.4), Inches(0.34),
        "20期频率差 · 四位置接力观察", 16, COLORS["gold"], True
    )
    fixed.add_text(
        cover, Inches(0.84), Inches(5.88), Inches(8.5), Inches(0.34),
        "挂机前规则说明｜结果等待实际运行", 14, COLORS["white"]
    )
    fixed.set_notes(
        cover,
        "本期只说明邻位热差迁移怎样计算、怎样运行和怎样记录，不提前给出结果。"
    )

    slide = body_slide(prs, "这个实验到底在问什么", "实验问题")
    card(
        slide, Inches(0.78), Inches(1.8), Inches(5.55), Inches(3.6),
        "左边热，右边冷",
        "同一个数字在相邻两个位置的出现次数可能不同。本实验只问：左侧位置偏热、右侧目标位置偏冷的数字，向右迁移后是否值得继续观察。",
        "gold", 18, 18
    )
    card(
        slide, Inches(6.58), Inches(1.8), Inches(5.55), Inches(3.6),
        "先冻结，再向前跑",
        f"用最近{LOOKBACK}期计算四组三码，冻结后轮投{RUN_PERIODS}期。真实命中、投入和波动只由实际挂机记录回答。",
        "blue", 18, 18
    )
    fixed.set_notes(slide, "强调这是一个可证伪问题，不把频率差解释成必然规律。")

    slide = body_slide(prs, "三枚号码怎样选出来", "核心规则")
    card(
        slide, Inches(0.78), Inches(1.8), Inches(3.55), Inches(3.95),
        "第一步｜数左边",
        f"统计最近{LOOKBACK}期，数字0到9在来源位置各出现多少次。",
        "gold", 17, 17
    )
    card(
        slide, Inches(4.65), Inches(1.8), Inches(3.55), Inches(3.95),
        "第二步｜减右边",
        "同一个数字的来源次数减去目标次数，得到“左热右冷”的频率差。",
        "blue", 17, 17
    )
    card(
        slide, Inches(8.52), Inches(1.8), Inches(3.55), Inches(3.95),
        "第三步｜取前三",
        "按频率差、来源次数和最近出现顺序排序，取前三个数字；冻结期间不再修改。",
        "green", 17, 17
    )
    fixed.set_notes(slide, "用三步讲清计算，不需要观众记公式，只要理解比较的是相邻位置差值。")

    slide = body_slide(prs, "本轮四组号码已经冻结", "冻结快照")
    for i, (route, play, _, _) in enumerate(ROUTES):
        x = Inches(0.72 + (i % 2) * 6.08)
        y = Inches(1.76 + (i // 2) * 2.18)
        accent = ["gold", "blue", "green", "red"][i]
        card(
            slide, x, y, Inches(5.78), Inches(1.82),
            f"{route}｜投{play}",
            " ".join(map(str, groups[i])),
            accent, 17, 26
        )
    fixed.add_text(
        slide, Inches(0.9), Inches(6.15), Inches(11.2), Inches(0.38),
        f"计算截止期：{issue}｜四组三码冻结{RUN_PERIODS}期",
        16, COLORS["white"], True, PP_ALIGN.CENTER
    )
    fixed.set_notes(slide, "逐组读出来源、目标和号码。提醒观众这些号码在本轮内固定。")

    slide = body_slide(prs, "四份主方案怎样轮流运行", "执行规则")
    labels = [
        f"万到千\n{' '.join(map(str, groups[0]))}",
        f"千到百\n{' '.join(map(str, groups[1]))}",
        f"百到十\n{' '.join(map(str, groups[2]))}",
        f"十到个\n{' '.join(map(str, groups[3]))}",
    ]
    for i, label in enumerate(labels):
        card(
            slide, Inches(0.52 + i * 3.17), Inches(2.02), Inches(2.88), Inches(2.0),
            f"第{i + 1}份", label, ["gold", "blue", "green", "red"][i], 15, 20
        )
        if i < 3:
            fixed.add_text(
                slide, Inches(3.38 + i * 3.17), Inches(2.72), Inches(0.35), Inches(0.4),
                "→", 22, COLORS["gray"], True, PP_ALIGN.CENTER
            )
    fixed.add_text(
        slide, Inches(0.9), Inches(4.62), Inches(11.1), Inches(0.5),
        "手工开启顶部“方案轮投”｜每期只运行一份｜连续24期",
        18, COLORS["white"], True, PP_ALIGN.CENTER
    )
    fixed.add_text(
        slide, Inches(1.0), Inches(5.5), Inches(10.9), Inches(0.48),
        "第一期实际从哪一份开始，以软件显示为准并记录。",
        18, COLORS["red"], True, PP_ALIGN.CENTER
    )
    fixed.set_notes(slide, "顶部轮投属于软件手动开关。运行期间不要重启，避免起点发生变化。")

    slide = body_slide(prs, "一轮从开始到结束怎么做", "操作流程")
    card(
        slide, Inches(0.78), Inches(1.72), Inches(3.45), Inches(4.35),
        "开始前",
        "导入四份主方案。核对位置和号码。开启顶部方案轮投。记录第一期实际起点。",
        "gold", 17, 16
    )
    card(
        slide, Inches(4.45), Inches(1.72), Inches(3.45), Inches(4.35),
        "运行中",
        "每期记录实际方案、目标位置、三枚号码、开奖号和是否命中。全程平倍，不改号、不重排。",
        "blue", 17, 16
    )
    card(
        slide, Inches(8.12), Inches(1.72), Inches(3.45), Inches(4.35),
        "第24期后",
        "立即停止主组并保存记录。随机对照另开一轮24期，不与主组同时并投。",
        "green", 17, 16
    )
    fixed.set_notes(slide, "这是完整操作示例，不是假设命中案例。")

    slide = body_slide(prs, "随机对照为什么必须分开", "对照方法")
    card(
        slide, Inches(0.78), Inches(1.82), Inches(5.55), Inches(3.5),
        "邻位热差主组",
        "每期一个位置、3个冻结号码。四个位置按实际轮投顺序接力，共运行24期。",
        "gold", 18, 18
    )
    card(
        slide, Inches(6.58), Inches(1.82), Inches(5.55), Inches(3.5),
        "随机三码对照",
        "位置、号码数量和运行期数保持一致，另行运行24期，只改变号码产生方式。",
        "blue", 18, 18
    )
    fixed.add_text(
        slide, Inches(1.0), Inches(5.62), Inches(10.9), Inches(0.52),
        "两组同时并投会翻倍成本，也会破坏比较口径。",
        20, COLORS["red"], True, PP_ALIGN.CENTER
    )
    fixed.set_notes(slide, "对照只改变选号方法，其他条件尽量保持一致。")

    slide = body_slide(prs, "准备多少资金，什么时候停", "资金与边界")
    metric(
        slide, Inches(0.82), Inches(1.9), Inches(3.55),
        "建议本金", "144元", "主组72元 + 对照72元", "gold"
    )
    metric(
        slide, Inches(4.89), Inches(1.9), Inches(3.55),
        "止盈", "不设置", "固定期数保证样本完整", "blue"
    )
    metric(
        slide, Inches(8.96), Inches(1.9), Inches(3.55),
        "止损", "不设置", "每组24期为硬边界", "green"
    )
    card(
        slide, Inches(0.82), Inches(3.54), Inches(11.7), Inches(2.05),
        "硬停止条件",
        "主组24期结束即停；随机对照24期结束即停。每个号码按1元计算，全程平倍，不追损，不在轮内重启。",
        "red", 18, 20
    )
    fixed.set_notes(slide, "144元是完成主组和对照的毛投入上限，不是盈利目标。")

    slide = body_slide(prs, "24期以后怎样判断", "判定方法")
    card(
        slide, Inches(0.78), Inches(1.78), Inches(5.55), Inches(3.7),
        "只看真实记录",
        "比较主组和随机对照的命中次数、连续未中、实际投入，以及不同起点是否造成明显差异。",
        "gold", 18, 18
    )
    card(
        slide, Inches(6.58), Inches(1.78), Inches(5.55), Inches(3.7),
        "不临时修规则",
        "不能删除差的路线、改变起点、加倍或延长观察期后再说原方案有效。任何改动都必须作为下一版重新验证。",
        "blue", 18, 18
    )
    fixed.add_text(
        slide, Inches(1.0), Inches(5.72), Inches(10.9), Inches(0.5),
        "先完成运行，再制作结果复盘。",
        20, COLORS["green"], True, PP_ALIGN.CENTER
    )
    fixed.set_notes(slide, "结尾只说明判定纪律，不提前宣布方案好坏。")

    fixed.add_end(prs, end_png)
    PPT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPT_PATH)


def ppt_text(prs):
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
        text.append(slide.notes_slide.notes_text_frame.text)
    return "\n".join(text)


def validate(groups, issue: str, audit):
    txts = list(PACKAGE.glob("*.txt"))
    if len(txts) != 8:
        raise AssertionError(f"TXT数量错误: {len(txts)}")
    for path in PACKAGE.iterdir():
        if re.search(r"(^|[_-])[BC]\d{3}([_-]|$)|SET_\d+", path.name, re.I):
            raise AssertionError(f"对外交付文件名存在工程编号: {path.name}")
    for path in txts:
        raw = path.read_bytes()
        if b"\r\n" not in raw or b"SchemeCreator=\r\n" not in raw:
            raise AssertionError(f"TXT编码或字段异常: {path.name}")

    prs = Presentation(PPT_PATH)
    text = ppt_text(prs)
    if len(prs.slides) != 10:
        raise AssertionError(f"PPT页数错误: {len(prs.slides)}")
    if any(slide._element.get("show") == "0" for slide in prs.slides):
        raise AssertionError("PPT仍存在隐藏页")
    for index, slide in enumerate(prs.slides, 1):
        if not slide.notes_slide.notes_text_frame.text.strip():
            raise AssertionError(f"第{index}页缺少备注")
    for token in [
        "B001", "B002", "B003", "B004", "C001", "C002", "C003", "C004",
        "B395", "SET_001", "方案ID", "批次ID",
    ]:
        if token in text:
            raise AssertionError(f"PPT出现工程编号: {token}")
    if re.search(r"\d+(?:\.\d+)?U\b", text):
        raise AssertionError("PPT金额仍使用U")
    for phrase in ["隐藏附录", "固定第二页", "已经盈利", "已经亏损", "优于随机", "验证失败"]:
        if phrase in text:
            raise AssertionError(f"PPT出现禁用内容: {phrase}")
    for phrase in ["144元", "实际挂机记录", "结果等待实际运行"]:
        if phrase not in text:
            raise AssertionError(f"PPT缺少关键内容: {phrase}")

    manifest = {
        "internal_scheme_id": INTERNAL_SCHEME_ID,
        "stage": "PRE_RUN_SETUP",
        "issue_cutoff": issue,
        "lookback": LOOKBACK,
        "run_periods": RUN_PERIODS,
        "groups": {
            ROUTES[i][0]: {
                "play": ROUTES[i][1],
                "digits": list(groups[i]),
            }
            for i in range(4)
        },
        "capital_yuan": 144,
        "take_profit": "NONE",
        "stop_loss": "NONE",
        "hard_stop": "主组24期；随机对照24期",
        "zip": ZIP_PATH.name,
        "ppt": PPT_PATH.name,
        "zip_sha256": sha256(ZIP_PATH),
        "ppt_sha256": sha256(PPT_PATH),
        "ppt_slides": len(prs.slides),
        "hidden_slides": 0,
        "currency": "元",
        "rolling_audit": audit,
        "status": "BUILD_VALIDATED_AWAITING_RUNTIME",
    }
    MANIFEST.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


def main():
    if OUT_ROOT.exists():
        shutil.rmtree(OUT_ROOT)
    OUT_ROOT.mkdir(parents=True)
    rows = parse_draws(INPUT_DATA)
    groups, _ = freeze_groups(rows)
    issue = rows[-1][0]
    audit = rolling_audit(rows)
    build_package(groups, issue)
    build_ppt(groups, issue)
    validate(groups, issue, audit)
    print(
        "DELIVERY_OK name=邻位热差迁移 stage=PRE_RUN_SETUP "
        "currency=元 periods=24 controls=RANDOM"
    )


if __name__ == "__main__":
    main()
