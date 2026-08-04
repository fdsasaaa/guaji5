#!/usr/bin/env python3
from __future__ import annotations
from collections import Counter
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile
import hashlib, json, re, shutil, sys
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT=Path(__file__).resolve().parents[1]
sys.path.insert(0,str(ROOT/'tools'))
import materialize_ppt_fixed_pages as fixed

INTERNAL_SCHEME_ID='B394-SET-001'
INPUT_DATA=ROOT/'01_本次输入'/'哈希分分彩_20260731_0181至0380.txt'
OUT_ROOT=ROOT/'dist'/'冷热三角快照接力_交付'
PACKAGE=OUT_ROOT/'package'
ZIP_PATH=OUT_ROOT/'冷热三角快照接力_方案套.zip'
PPT_PATH=OUT_ROOT/'冷热三角快照接力_挂机前讲解.pptx'
MANIFEST=OUT_ROOT/'DELIVERY_MANIFEST.json'
COLORS={'bg':(8,12,17),'panel':(18,25,32),'white':(244,247,250),'gray':(166,178,188),'gold':(232,177,64),'green':(65,201,151),'red':(226,93,93),'blue':(77,153,230),'line':(68,82,96)}


def sha256(path:Path)->str:
 h=hashlib.sha256()
 with path.open('rb') as f:
  for c in iter(lambda:f.read(1024*1024),b''): h.update(c)
 return h.hexdigest()

def parse_draws(path:Path):
 rows=[]
 for raw in path.read_text(encoding='utf-8').splitlines():
  if not raw.strip(): continue
  issue,num=raw.strip().split('=',1)
  if not issue.isdigit() or len(num)!=5 or not num.isdigit(): raise ValueError(f'开奖格式错误: {raw}')
  rows.append((issue,[int(x) for x in num]))
 if len(rows)<18: raise ValueError('开奖数据不足18期')
 return rows

def rank_freq(values):
 c=Counter(values); last={d:max((i for i,x in enumerate(values) if x==d),default=-1) for d in range(10)}
 return sorted(range(10),key=lambda d:(-c[d],-last[d],d))
def omission3(values):
 om={}
 for d in range(10):
  n=0
  for x in reversed(values[-18:]):
   if x==d: break
   n+=1
  om[d]=n
 return tuple(sorted(sorted(range(10),key=lambda d:(-om[d],d))[:3]))
def snapshot(rows):
 return [tuple(sorted(rank_freq([d[2] for _,d in rows[-6:]])[:3])),tuple(sorted(rank_freq([d[3] for _,d in rows[-12:]])[3:6])),omission3([d[4] for _,d in rows])]

def common(strategy,play,enabled):
 return ['True' if enabled else 'False',strategy,'软件名称=CXGGJ','玩法类型=定位胆',f'玩法名称={play}','金额模式=2','投注监控=False-','投注监控模式=0','任选中奖=1-10','任选位置=']
def tail():
 return ['翻倍方式=0','正集=True','倍投类型=0','倍投计划=1,1,1,1,1,1,1,1,1,1','倍投方案=1,1,1,1,1,1,1,1,1,1','显示更多=False','真实投注1=False-50000','真实投注2=False-50000','模拟投注1=False-50000','模拟投注2=False-50000','盈利跳转=False-50000-1','亏损跳转=False-50000-1','盈利停止=False-50000','亏损停止=False-50000','投注时间=False','投注时间类型=0','范围开始时间=False-09:01:00','范围停止时间=False-21:32:00','范围停止类型=0','倒计时停止时间=02:00:00','倒计时停止类型=0']
def fixed_txt(play,digits,enabled=True):
 lines=common('定码轮换',play,enabled)+['换号规则=9','换号期数=18']+tail()+[f"定码轮换内容={' '.join(map(str,digits))}",'定码轮换单组=True','SchemeCreator=']
 return '\r\n'.join(lines)+'\r\n'
def random_txt(play):
 lines=common('随机出号',play,False)+['换号规则=10','换号期数=1']+tail()+['随机出号模板=模板1','随机出号个数=3','SchemeCreator=']
 return '\r\n'.join(lines)+'\r\n'

def build_package(groups,issue):
 if PACKAGE.exists(): shutil.rmtree(PACKAGE)
 PACKAGE.mkdir(parents=True)
 internal=[('短热百位-定码轮换.txt','百位',groups[0]),('中温十位-定码轮换.txt','十位',groups[1]),('长遗漏个位-定码轮换.txt','个位',groups[2])]
 for name,play,digits in internal: (PACKAGE/name).write_bytes(fixed_txt(play,digits).encode('gbk'))
 for name,play in [('随机百位-随机出号.txt','百位'),('随机十位-随机出号.txt','十位'),('随机个位-随机出号.txt','个位')]: (PACKAGE/name).write_bytes(random_txt(play).encode('gbk'))
 readme=f'''# 冷热三角快照接力｜挂机前说明

本方案处于“挂机前验证准备”阶段。本文件只说明规则和运行方法，不提前判断盈利、亏损或方案好坏。

## 本轮三组号码

- 短热百位：{' '.join(map(str,groups[0]))}
- 中温十位：{' '.join(map(str,groups[1]))}
- 长遗漏个位：{' '.join(map(str,groups[2]))}
- 计算截止期：{issue}

## 运行方式

导入三份主方案后，手工勾选软件顶部“方案轮投”。每期只运行一份，连续观察18期。随机对照必须单独运行18期，不能与主组同时并投。

## 资金边界

- 每个号码按1元计算；每期3个号码，共3元。
- 主组18期毛投入54元；随机对照18期毛投入54元。
- 建议准备108元。
- 止盈：不设置。
- 止损：不设置。
- 替代停止条件：每组18期结束立即停止，不临时改码、不追加倍投。

## 需要记录

记录每期实际运行的自然方案名称、位置、号码、开奖号和是否命中。完成实际挂机后，再根据真实记录制作结果复盘PPT。
'''
 (PACKAGE/'00_使用说明.md').write_text(readme,encoding='utf-8')
 checklist='''# 导入与运行核对表

- [ ] 三份主方案和三份随机对照均可导入
- [ ] 主组默认勾选，对照默认不勾选
- [ ] 顶部“方案轮投”已手工开启
- [ ] 每期只运行一份方案
- [ ] 每个号码按1元计算，全程平倍
- [ ] 主组18期结束后停止
- [ ] 随机对照另行运行18期
- [ ] 已记录每期实际方案、开奖号和命中情况
- [ ] 未在轮内修改号码、顺序或倍投
'''
 (PACKAGE/'01_导入运行核对表.md').write_text(checklist,encoding='utf-8')
 record='''期号,实际方案名称,位置,号码,开奖号,是否命中,备注
'''
 (PACKAGE/'02_实际挂机记录模板.csv').write_text(record,encoding='utf-8-sig')
 if ZIP_PATH.exists(): ZIP_PATH.unlink()
 with ZipFile(ZIP_PATH,'w',ZIP_DEFLATED) as z:
  for p in sorted(PACKAGE.iterdir()): z.write(p,arcname=p.name)

def body_slide(prs,title,kicker=''):
 s=prs.slides.add_slide(prs.slide_layouts[6]); f=s.background.fill; f.solid(); f.fore_color.rgb=RGBColor(*COLORS['bg'])
 if kicker: fixed.add_text(s,Inches(.78),Inches(.38),Inches(4),Inches(.3),kicker,12,COLORS['gold'],True)
 fixed.add_text(s,Inches(.78),Inches(.72),Inches(11.7),Inches(.58),title,28,COLORS['white'],True)
 line=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,Inches(.78),Inches(1.42),Inches(11.74),Inches(.025)); line.fill.solid(); line.fill.fore_color.rgb=RGBColor(*COLORS['line']); line.line.fill.background()
 return s
def card(s,x,y,w,h,title,body,accent='gold',ts=17,bs=16):
 sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,x,y,w,h); sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor(*COLORS['panel']); sh.line.color.rgb=RGBColor(*COLORS['line'])
 fixed.add_text(s,x+Inches(.24),y+Inches(.16),w-Inches(.48),Inches(.36),title,ts,COLORS[accent],True)
 fixed.add_text(s,x+Inches(.24),y+Inches(.58),w-Inches(.48),h-Inches(.76),body,bs,COLORS['white'],False,PP_ALIGN.LEFT,MSO_ANCHOR.TOP)
def metric(s,x,y,w,label,value,sub='',accent='gold'):
 sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,x,y,w,Inches(1.15)); sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor(*COLORS['panel']); sh.line.color.rgb=RGBColor(*COLORS['line'])
 fixed.add_text(s,x+Inches(.18),y+Inches(.12),w-Inches(.36),Inches(.22),label,11,COLORS['gray'])
 fixed.add_text(s,x+Inches(.18),y+Inches(.36),w-Inches(.36),Inches(.42),value,24,COLORS[accent],True)
 if sub: fixed.add_text(s,x+Inches(.18),y+Inches(.82),w-Inches(.36),Inches(.2),sub,10,COLORS['gray'])

def build_ppt(groups,issue):
 cover_png=ROOT/'assets'/'ppt'/'fixed_pages'/'首页背景图谱.png'; end_png=ROOT/'assets'/'ppt'/'fixed_pages'/'固定最后一页_画面.png'
 prs=fixed.new_prs(); cover=fixed.add_cover(prs,cover_png)
 fixed.add_text(cover,Inches(.82),Inches(1.05),Inches(8.8),Inches(.62),'冷热三角接力',32,COLORS['white'],True)
 fixed.add_text(cover,Inches(.84),Inches(1.82),Inches(8.4),Inches(.34),'18期快照 · 三位置轮流观察',16,COLORS['gold'],True)
 fixed.add_text(cover,Inches(.84),Inches(5.88),Inches(8.5),Inches(.34),'挂机前规则说明｜结果等待实际运行',14,COLORS['white'])
 fixed.set_notes(cover,'本期只介绍怎样运行冷热三角接力，不提前判断它是否盈利或有效。')

 s=body_slide(prs,'为什么做这个实验','实验问题')
 card(s,Inches(.78),Inches(1.8),Inches(5.55),Inches(3.6),'三个观察角度','短热、中温、长遗漏分别代表不同时间窗口。把它们放在百位、十位和个位轮流运行，观察结构是否具有稳定表现。','gold',18,18)
 card(s,Inches(6.58),Inches(1.8),Inches(5.55),Inches(3.6),'先问问题，不给答案','这次先冻结规则并向前运行18期。实际命中、投入和波动都由挂机记录决定，不用历史数据替代最终答案。','blue',18,18)
 fixed.set_notes(s,'强调这是待验证问题。观众现在只需要知道为什么设计这套结构。')

 s=body_slide(prs,'三组号码怎样产生','核心规则')
 for i,(title,body,accent) in enumerate([
  ('短热百位',f"最近6期百位按出现次数排序，取前3个。\n本轮：{' '.join(map(str,groups[0]))}",'gold'),
  ('中温十位',f"最近12期十位排除最热前三，再取第4至第6名。\n本轮：{' '.join(map(str,groups[1]))}",'blue'),
  ('长遗漏个位',f"最近18期个位按当前遗漏从高到低取3个。\n本轮：{' '.join(map(str,groups[2]))}",'green')]):
  card(s,Inches(.72+i*4.04),Inches(1.82),Inches(3.82),Inches(3.85),title,body,accent,18,17)
 fixed.add_text(s,Inches(.9),Inches(5.95),Inches(11.2),Inches(.42),f'计算截止期：{issue}｜号码冻结18期，轮内不修改',18,COLORS['white'],True,PP_ALIGN.CENTER)
 fixed.set_notes(s,'逐一说明三个窗口。号码在18期内固定，避免边跑边改规则。')

 s=body_slide(prs,'三份主方案怎样轮流运行','执行规则')
 for i,(title,body,accent) in enumerate([
  ('短热百位',f"号码 {' '.join(map(str,groups[0]))}",'gold'),('中温十位',f"号码 {' '.join(map(str,groups[1]))}",'blue'),('长遗漏个位',f"号码 {' '.join(map(str,groups[2]))}",'green')]):
  card(s,Inches(.88+i*4.05),Inches(2.05),Inches(3.55),Inches(2.02),title,body,accent,17,22)
  if i<2: fixed.add_text(s,Inches(4.42+i*4.05),Inches(2.72),Inches(.42),Inches(.45),'→',24,COLORS['gray'],True,PP_ALIGN.CENTER)
 fixed.add_text(s,Inches(.95),Inches(4.62),Inches(11),Inches(.56),'手工开启顶部“方案轮投”｜每期只运行一份｜18期后停止',18,COLORS['white'],True,PP_ALIGN.CENTER)
 fixed.add_text(s,Inches(1.2),Inches(5.52),Inches(10.5),Inches(.5),'第一期实际运行哪一份，以软件显示为准并记录。',18,COLORS['red'],True,PP_ALIGN.CENTER)
 fixed.set_notes(s,'PPT使用自然名称，不展示内部文件编号。顶部方案轮投需要人工开启。')

 s=body_slide(prs,'一轮怎样从开始跑到结束','演示例子')
 card(s,Inches(.78),Inches(1.72),Inches(3.45),Inches(4.35),'开始前','导入三份主方案。确认号码和位置。开启顶部方案轮投。记录第一期实际起点。','gold',17,16)
 card(s,Inches(4.45),Inches(1.72),Inches(3.45),Inches(4.35),'运行中','每期记录自然方案名称、位置、号码、开奖号和是否命中。轮内不改码、不改顺序、不加倍。','blue',17,16)
 card(s,Inches(8.12),Inches(1.72),Inches(3.45),Inches(4.35),'第18期后','立即停止本轮。保存完整记录。重新计算下一张快照，旧号码不继续使用。','green',17,16)
 fixed.set_notes(s,'这是操作演示，不是命中案例。重点是把完整运行过程讲清。')

 s=body_slide(prs,'随机对照必须单独运行','对照方法')
 card(s,Inches(.78),Inches(1.82),Inches(5.55),Inches(3.5),'冷热三角主组','每期一个位置、3个号码。连续运行18期，号码在轮前冻结。总毛投入54元。','gold',18,18)
 card(s,Inches(6.58),Inches(1.82),Inches(5.55),Inches(3.5),'随机三码对照','位置和号码数量完全相同，另行运行18期。总毛投入同样为54元。','blue',18,18)
 fixed.add_text(s,Inches(1),Inches(5.62),Inches(10.9),Inches(.52),'两组不能同时并投，否则成本和结果都无法比较。',20,COLORS['red'],True,PP_ALIGN.CENTER)
 fixed.set_notes(s,'对照组与主组分开运行，保证比较口径一致。')

 s=body_slide(prs,'准备多少资金，什么时候停止','资金与边界')
 metric(s,Inches(.82),Inches(1.9),Inches(3.55),'建议本金','108元','主组54元 + 对照54元','gold')
 metric(s,Inches(4.89),Inches(1.9),Inches(3.55),'止盈','不设置','固定期数保证样本完整','blue')
 metric(s,Inches(8.96),Inches(1.9),Inches(3.55),'止损','不设置','每组18期为硬边界','green')
 card(s,Inches(.82),Inches(3.54),Inches(11.7),Inches(2.05),'硬停止条件','主组18期结束即停；随机对照18期结束即停。每个号码按1元计算，全程平倍，不追损，不在轮内重启。','red',18,20)
 fixed.set_notes(s,'108元是两组完整观察的毛投入上限，不是盈利目标。')

 s=body_slide(prs,'挂机时必须记录什么','记录要求')
 card(s,Inches(.78),Inches(1.78),Inches(5.55),Inches(3.7),'每期记录','期号、实际方案名称、位置、三枚号码、开奖号、是否命中，以及软件是否发生重启。','gold',18,18)
 card(s,Inches(6.58),Inches(1.78),Inches(5.55),Inches(3.7),'保持不变','18期内不改号码、不换起点、不调整倍数。任何临时修改都会让本轮验证失去统一口径。','blue',18,18)
 fixed.add_text(s,Inches(1),Inches(5.72),Inches(10.9),Inches(.5),'真实记录完成后，再制作结果复盘PPT。',20,COLORS['green'],True,PP_ALIGN.CENTER)
 fixed.set_notes(s,'这里告诉观众结果从哪里来：来自实际挂机记录，不是提前算出的答案。')

 s=body_slide(prs,'18期结束后再回答好不好','等待验证')
 card(s,Inches(.92),Inches(1.82),Inches(11.45),Inches(3.55),'本期只冻结验证方法','现在可以确认的是规则、号码、资金和记录方式。是否盈利、是否优于随机、哪个起点更稳定，都要等实际挂机完成后再判断。','green',19,21)
 fixed.add_text(s,Inches(1.1),Inches(5.62),Inches(10.7),Inches(.55),'先运行，后复盘；不让历史数据替代真实结果。',20,COLORS['white'],True,PP_ALIGN.CENTER)
 fixed.set_notes(s,'主体收尾不下结果结论，只说明下一步是完成真实运行。')
 fixed.add_end(prs,end_png)
 PPT_PATH.parent.mkdir(parents=True,exist_ok=True); prs.save(PPT_PATH)

def ppt_text(prs):
 out=[]
 for s in prs.slides:
  for sh in s.shapes:
   if hasattr(sh,'text'): out.append(sh.text)
  out.append(s.notes_slide.notes_text_frame.text)
 return '\n'.join(out)
def validate():
 txts=list(PACKAGE.glob('*.txt'))
 if len(txts)!=6: raise AssertionError(f'TXT数量错误: {len(txts)}')
 for p in PACKAGE.iterdir():
  if re.search(r'(^|[_-])[BC]\d{3}([_-]|$)|SET_\d+',p.name,re.I): raise AssertionError(f'对外交付文件名存在工程编号: {p.name}')
 for p in txts:
  raw=p.read_bytes()
  if b'\r\n' not in raw or b'SchemeCreator=\r\n' not in raw: raise AssertionError(f'TXT编码或字段异常: {p.name}')
 prs=Presentation(PPT_PATH); text=ppt_text(prs)
 if len(prs.slides)!=10: raise AssertionError(f'PPT页数错误: {len(prs.slides)}')
 if any(s._element.get('show')=='0' for s in prs.slides): raise AssertionError('PPT仍存在隐藏页')
 for i,s in enumerate(prs.slides,1):
  if not s.notes_slide.notes_text_frame.text.strip(): raise AssertionError(f'第{i}页缺少备注')
 for token in ['B001','B002','B003','C001','C002','C003','B394','SET_001','方案ID','批次ID']:
  if token in text: raise AssertionError(f'PPT出现工程编号: {token}')
 if re.search(r'\d+(?:\.\d+)?U\b',text): raise AssertionError('PPT金额仍使用U')
 for phrase in ['固定第二页','先选平台，再谈方案','隐藏附录','优势未证实','证明不盈利']:
  if phrase in text: raise AssertionError(f'PPT出现禁用内容: {phrase}')
 for phrase in ['108元','实际挂机记录','结果等待实际运行']:
  if phrase not in text: raise AssertionError(f'PPT缺少关键内容: {phrase}')
 data={'internal_scheme_id':INTERNAL_SCHEME_ID,'zip':ZIP_PATH.name,'ppt':PPT_PATH.name,'zip_sha256':sha256(ZIP_PATH),'ppt_sha256':sha256(PPT_PATH),'ppt_slides':len(prs.slides),'hidden_slides':0,'currency':'元','stage':'PRE_RUN_SETUP','fixed_second_page':False,'status':'BUILD_VALIDATED_AWAITING_RUNTIME'}
 MANIFEST.write_text(json.dumps(data,ensure_ascii=False,indent=2)+'\n',encoding='utf-8')

def main():
 if OUT_ROOT.exists(): shutil.rmtree(OUT_ROOT)
 OUT_ROOT.mkdir(parents=True)
 rows=parse_draws(INPUT_DATA); groups=snapshot(rows); issue=rows[-1][0]
 build_package(groups,issue); build_ppt(groups,issue); validate()
 print('DELIVERY_OK name=冷热三角快照接力 stage=PRE_RUN_SETUP currency=元 second=DISABLED appendix=DISABLED')
if __name__=='__main__': main()
